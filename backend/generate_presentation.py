from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

def create_executive_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    DARK_BLUE = RGBColor(15, 23, 42)
    ACCENT_BLUE = RGBColor(59, 130, 246)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(148, 163, 184)
    
    def add_title_shape(slide, text, top=0.3, font_size=40):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_BLUE
        shape.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.color.rgb = WHITE
        return shape
    
    def add_content_box(slide, left, top, width, height, title, bullets):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(30, 41, 59)
        box.line.color.rgb = ACCENT_BLUE
        box.line.width = Pt(2)
        
        title_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.15), Inches(width - 0.4), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        
        content_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.6), Inches(width - 0.4), Inches(height - 0.8))
        tf = content_box.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE
            p.space_after = Pt(8)
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()
    
    title = slide1.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Enterprise GRC Platform"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    subtitle = slide1.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12), Inches(1))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "Unified Governance, Risk & Compliance Management"
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    tagline = slide1.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(0.8))
    tf = tagline.text_frame
    p = tf.paragraphs[0]
    p.text = "Multi-Tenant | AI-Powered | Enterprise-Grade Security"
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Platform Overview
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = DARK_BLUE
    bg2.line.fill.background()
    add_title_shape(slide2, "Platform Overview")
    
    add_content_box(slide2, 0.3, 1.7, 4, 2.5, "Single Source of Truth",
        ["Centralized GRC data repository", "Real-time risk visibility", "Cross-framework compliance tracking", "Unified audit trail"])
    
    add_content_box(slide2, 4.6, 1.7, 4, 2.5, "Multi-Tenancy",
        ["Complete tenant isolation", "Row-level security", "Role-based access control", "Enterprise SSO support"])
    
    add_content_box(slide2, 8.9, 1.7, 4, 2.5, "AI-Powered Intelligence",
        ["Automated control mapping", "Evidence quality assessment", "Gap analysis & recommendations", "Regulatory change impact analysis"])
    
    overview_text = slide2.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(2.5))
    tf = overview_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "The Enterprise GRC Platform provides a comprehensive solution for managing governance, risk, and compliance across your organization. Built with enterprise-grade security and multi-tenant architecture, it enables organizations to streamline GRC processes while maintaining complete data isolation and regulatory compliance."
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    
    # Slide 3: Core Modules
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = DARK_BLUE
    bg3.line.fill.background()
    add_title_shape(slide3, "Core Modules")
    
    modules = [
        ("Compliance Management", ["Multi-framework support (8+ frameworks)", "User-uploaded regulatory frameworks", "Control mapping & gap analysis", "Evidence management & AI assessment"]),
        ("Enterprise Risk Management", ["Comprehensive risk register", "Risk appetite management", "KRI tracking & monitoring", "Incident management & response"]),
        ("Governance & Policy", ["Policy lifecycle management", "Document version control", "Approval workflows", "Committee & board management"]),
        ("Regulatory Change Mgmt", ["Regulatory update tracking", "Impact assessment tools", "Implementation task tracking", "AI-powered gap analysis"]),
        ("Attestation & Certification", ["SOX 302/404 certifications", "Campaign-based attestations", "Cascade reminder workflows", "Policy sign-off tracking"]),
        ("Asset & Vulnerability", ["IT asset inventory", "Vulnerability management", "SLA tracking & escalation", "AI fix recommendations"])
    ]
    
    for i, (title, bullets) in enumerate(modules):
        col = i % 3
        row = i // 3
        left = 0.3 + col * 4.3
        top = 1.7 + row * 2.8
        add_content_box(slide3, left, top, 4, 2.6, title, bullets)
    
    # Slide 4: Key Features & Benefits
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    bg4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg4.fill.solid()
    bg4.fill.fore_color.rgb = DARK_BLUE
    bg4.line.fill.background()
    add_title_shape(slide4, "Key Features & Benefits")
    
    add_content_box(slide4, 0.3, 1.7, 6.3, 2.7, "Key Features",
        ["Unified Control Library with cross-framework mapping", "AI-powered evidence assessment & recommendations", "Real-time compliance dashboards & analytics", "Automated workflow & approval management", "Board & committee governance tracking", "Bulk import via CSV/Excel templates"])
    
    add_content_box(slide4, 6.8, 1.7, 6.3, 2.7, "Business Benefits",
        ["60% reduction in audit preparation time", "Single pane of glass for all GRC activities", "Reduced compliance costs through automation", "Improved risk visibility & decision making", "Regulatory change readiness", "Enhanced board reporting capabilities"])
    
    add_content_box(slide4, 0.3, 4.6, 12.7, 2.5, "Technical Advantages",
        ["Enterprise-grade multi-tenant architecture with complete data isolation",
         "Modern tech stack: FastAPI, Next.js 14, PostgreSQL with row-level security",
         "RESTful APIs for seamless integration with existing enterprise systems",
         "AI/ML capabilities powered by GPT-4 for intelligent automation",
         "Deterministic AI assessments ensuring reproducible compliance results"])
    
    # Slide 5: Implementation & Next Steps
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    bg5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg5.fill.solid()
    bg5.fill.fore_color.rgb = DARK_BLUE
    bg5.line.fill.background()
    add_title_shape(slide5, "Implementation & Next Steps")
    
    add_content_box(slide5, 0.3, 1.7, 4, 3, "Implementation Phases",
        ["Phase 1: Platform setup & configuration", "Phase 2: Framework upload & mapping", "Phase 3: Data migration & integration", "Phase 4: User training & rollout", "Phase 5: Ongoing support & optimization"])
    
    add_content_box(slide5, 4.6, 1.7, 4, 3, "Success Metrics",
        ["Compliance coverage rate", "Audit finding reduction", "Time-to-remediation", "Policy attestation rates", "Risk assessment completion"])
    
    add_content_box(slide5, 8.9, 1.7, 4, 3, "Support & Training",
        ["Dedicated implementation team", "Role-based training programs", "24/7 technical support", "Regular platform updates", "Best practice guidance"])
    
    contact = slide5.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12), Inches(1.5))
    tf = contact.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Ready to transform your GRC operations?"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "Contact us for a personalized demo and implementation roadmap"
    p2.font.size = Pt(18)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    
    output_path = "GRC_Platform_Executive_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_executive_presentation()
