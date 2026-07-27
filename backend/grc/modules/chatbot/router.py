"""
ComplyChatRouter - AI-Powered Compliance Q&A Integration
Integrates RAG-based chatbot into the main GRC platform
"""

from ...config import get_openai_model
import os
import sys
import json
import logging
import hashlib
import re
from io import BytesIO
from typing import List, Dict, Any, Optional, Sequence, Tuple
from pathlib import Path
from collections import deque
from datetime import datetime
from uuid import uuid4

from fastapi import APIRouter, Depends, HTTPException, status, Response, UploadFile, File, Form
from pydantic import AliasChoices, BaseModel, ConfigDict, Field, field_validator
from sqlalchemy.orm import Session
from sqlalchemy import text, func, or_

from ...services.ai_usage import usage_scope

# Add complychat to path
complychat_path = Path(__file__).parent / "complychat" / "complychat"
sys.path.insert(0, str(complychat_path))

logger = logging.getLogger(__name__)

# ============================================================================
# CONVERSATION HISTORY STORAGE
# ============================================================================
conversation_history: Dict[str, deque] = {}  # Backward-compatible history snapshots
session_message_logs: Dict[str, List[Dict[str, Any]]] = {}
session_uploaded_files: Dict[str, List[Dict[str, Any]]] = {}
langchain_histories: Dict[str, Any] = {}
MAX_HISTORY_LENGTH = 10
CHAT_UPLOAD_ROOT = Path(__file__).resolve().parents[3] / "uploads" / "complychat"
CHAT_UPLOAD_ROOT.mkdir(parents=True, exist_ok=True)
VECTOR_INDEX_STATE: Dict[int, Dict[str, str]] = {}
VECTOR_MATCH_LIMIT = int(os.getenv("COMPLYCHAT_VECTOR_MATCH_LIMIT", "16"))
VECTOR_SCORE_THRESHOLD = float(os.getenv("COMPLYCHAT_VECTOR_SCORE_THRESHOLD", "0.18"))
VECTOR_CATALOG_FALLBACK_LIMIT = int(os.getenv("COMPLYCHAT_VECTOR_CATALOG_FALLBACK_LIMIT", "120"))
VECTOR_DOC_LIMIT_ON_ASK = int(os.getenv("COMPLYCHAT_VECTOR_DOC_LIMIT_ON_ASK", "1500"))
VECTOR_PARSED_CONTROL_LIMIT = int(os.getenv("COMPLYCHAT_VECTOR_PARSED_CONTROL_LIMIT", "20000"))
VECTOR_SEEDED_CONTROL_LIMIT = int(os.getenv("COMPLYCHAT_VECTOR_SEEDED_CONTROL_LIMIT", "20000"))
VECTOR_POLICY_STATEMENT_LIMIT = int(os.getenv("COMPLYCHAT_VECTOR_POLICY_STATEMENT_LIMIT", "4000"))
VECTOR_GOV_VERSION_LIMIT = int(os.getenv("COMPLYCHAT_VECTOR_GOV_VERSION_LIMIT", "1200"))
VECTOR_ROUTE_TERMS = (
    "document", "policy", "charter", "framework", "control", "clause", "section",
    "requirement", "obligation", "assessment", "evidence", "uploaded", "excel", "sheet",
    "pdf", "docx", "what does", "what says", "based on", "from",
    "compare", "comparison", "align", "alignment", "gap", "difference",
)
VECTOR_SEMANTIC_TERMS = (
    "what does", "what says", "explain", "summarize", "summarise", "clause",
    "section", "requirement", "obligation", "policy says", "document says",
    "compare", "comparison", "align", "alignment", "gap", "difference",
)
# IMPORTANT: Only terms that have NO database tables (truly removed UI features).
# Audit findings, plans, engagements, reports ARE in the DB — route those to database instead.
# This list is intentionally narrow to avoid blocking valid DB queries.
AUDIT_QUERY_TERMS: tuple = ()  # Audit Management tables NOW exist — all route to database
GRC_RELEVANT_TERMS = (
    # Platform modules
    "grc", "erm", "enterprise risk", "risk management", "risk register",
    # Core GRC concepts
    "compliance", "framework", "control", "policy", "risk", "evidence", "governance",
    "regulatory", "regulation", "obligation", "maturity", "readiness",
    # Security / vulnerability
    "vulnerabilit", "vulnerabilty", "vunerabilit",  # typo-tolerant stems
    "pentest", "penetration", "cve", "patch", "exploit", "remediat", "threat", "attack",
    "cybersecurit", "cyber secur", "infosec", "security",
    # Asset management
    "asset", "inventor", "hardware", "software", "system", "infrastructure",
    # Vendor / third-party
    "vendor", "third party", "third-party", "supplier", "outsourc",
    # Incidents / exceptions
    "incident", "exception", "breach", "finding",
    # Certifications / journeys
    "attestation", "certification", "journey", "assessment", "gap", "audit",
    # Frameworks (abbreviations and names)
    "iso", "pci", "nist", "sbp", "sama", "dora", "gdpr", "soc 2", "soc2",
    "cobit", "cis", "hipaa", "sox", "basel",
    # Platform sub-modules & pages
    "rcsa", "self-assessment", "self assessment",
    "internal control", "key control", "control test",
    "regulatory change", "regulatory update", "new regulation",
    "oversight action", "committee meeting", "board meeting",
    "policy statement", "policy gap", "gap analysis",
    "pentest report", "vulnerability report", "vuln report",
    "control library", "control mapping", "control framework",
    "certification journey", "compliance program", "compliance status",
    "issue tracker", "open issue", "risk incident",
    "it asset", "asset inventory", "cde environment",
    "vendor assessment", "vendor review", "vendor risk",
    "policy exception", "control exception",
    # Org / people
    "organization", "tenant", "company", "committee", "document", "charter",
    "department", "business unit", "user", "role", "permission",
    # GRC status words
    "open", "overdue", "critical", "high risk", "pending", "active", "closed",
    # Audit Management (all tables in DB)
    "audit finding", "audit findings", "audit plan", "audit plans",
    "audit engagement", "audit engagements", "audit report", "audit reports",
    "audit workpaper", "workpaper", "pbc list", "prepared by client",
    "qaip", "audit quality", "audit maturity", "audit recommendation",
    "audit universe", "auditable entity", "audit team",
    "audit board pack", "board pack", "audit follow up",
    # CCM
    "ccm", "continuous control monitoring", "ccm rule", "ccm anomaly",
    "ccm exception", "control anomaly", "automated control",
    # Risk analytics
    "risk appetite", "risk tolerance", "likelihood impact", "risk matrix",
    "risk score history", "risk trend", "kri", "key risk indicator",
    "risk appetite config", "risk scoring", "risk appetite level",
    # Vulnerability sub-items
    "vulnerability mitigation", "vulnerable mitigation", "vuln mitigation",
    "vulnerability retest", "vulnerability sla", "sla breach", "sla target",
    "scan record", "scanner sync", "sync history", "integration exception",
    # Meeting details
    "meeting agenda", "meeting minutes", "agenda item", "action points",
    # Policy gap details
    "gap finding", "gap findings", "policy gap finding", "gap analysis result",
    "gap analysis results", "policy gap", "compliance gap finding",
    # RCSA sub-items
    "rcsa finding", "rcsa response", "rcsa campaign",
    # Regulatory details
    "regulatory feed", "regulatory impact", "regulatory task",
    # Common typos / misspellings
    "complaince", "complianse", "complience",  # compliance
    "goveranc", "governanc", "govrnance",  # governance
    "framwork", "frameowrk",  # framework
    "attesstati", "attestaion",  # attestation
    "regultory", "regualtory",  # regulatory
    "vulnerablity", "vunlerabilit",  # vulnerability
    "certifcation", "certifiation",  # certification
    "asessment", "assessement",  # assessment
    "controll", "cotrnol",  # control
    "incidnet",  # incident
)
OFF_TOPIC_TERMS = (
    "weather", "football", "cricket score", "movie review", "movie ticket",
    "recipe", "song lyrics", "joke", "dating", "bitcoin price", "crypto price",
    "celebrity", "travel booking", "restaurant", "shopping cart", "video game",
    "horoscope", "sports score", "stock price",
)
GREETING_TERMS = (
    "hi", "hello", "hey", "good morning", "good afternoon", "good evening", "thanks", "thank you"
)
FILE_ANALYSIS_TERMS = (
    "analyze", "analyse", "review", "summarize", "summary", "gap", "assess", "assessment",
    "what does this file", "check this document", "uploaded file", "uploaded files", "document",
    "comply", "complian", "this policy", "attached policy", "this document", "attached document",
    "extract controls", "identify gaps", "check gaps", "is there", "does this", "cover",
    "requirement", "obligation", "missing clause", "non-conformit",
)
DB_QUERY_TERMS = (
    # Explicit retrieval commands
    "show", "list", "count", "how many", "which", "what are", "find", "get", "fetch",
    "give me", "tell me", "display", "report on", "pull",
    # Natural question patterns about platform data
    "what about", "tell me about", "what is the", "what is our", "what are our",
    "details on", "details about", "information on", "info on",
    "how many", "how much", "total number", "count of", "number of",
    "give details", "provide details", "describe our", "describe the",
    # Status / state (always live-data signals)
    "status", "current status", "current state", "progress", "how is", "how are",
    "open", "closed", "resolved", "unresolved", "approved", "rejected", "in progress",
    "overdue", "missing", "weak", "critical", "top", "active", "pending", "on hold",
    # Quantity / existence checks
    "any open", "any pending", "any unresolved", "any active", "any critical",
    "are there", "is there any", "do we have", "how many open", "how many pending",
    "what is in", "what is the status", "what is currently",
    # Time-based = always live data
    "current", "recent", "latest", "last week", "this month", "today", "this year",
    # Summary / overview of live data
    "summary of", "overview of", "breakdown of", "breakdown", "distribution",
    "by severity", "by status", "by category", "by framework", "by department",
    # Possessive patterns = platform data
    "my risks", "my controls", "my vendors", "my exceptions", "our risks",
    "in the system", "in the platform", "in our platform", "in your system",
    "linked", "linked to", "associated with", "related to",
)
# Platform module nouns — if present, the question is about live DB data
# Covers FULL platform: ALL 20+ domains, all sub-modules, all pages, all features
# Banking / InfoSec specific terms included for production-grade coverage
PLATFORM_DATA_NOUNS = (
    # ── Risk Management ──
    "risk", "risks", "risk register", "risk incident", "risk incidents",
    "risk exception", "risk exceptions", "risk kri", "risk kris", "kri", "kris",
    "key risk indicator", "key risk indicators",
    "risk treatment", "risk treatments", "risk owner", "risk category",
    "risk review", "residual risk", "inherent risk", "risk appetite", "risk tolerance",
    "risk mitigation", "risk mitigation action", "risk mitigation actions",
    "risk score", "risk scoring", "risk score history", "risk trend",
    "risk assessment", "risk assessments", "risk assessment campaign",
    "risk report", "risk heat map", "risk heatmap",
    "likelihood", "impact scale", "risk matrix",
    # ── Exceptions & Issues ──
    "open exceptions", "policy exceptions", "exceptions", "control exceptions",
    "open issues", "issues", "issue tracker",
    # ── Incidents ──
    "open incidents", "security incidents", "incidents", "risk incidents", "incident log",
    "incident response",
    # ── Vendors ──
    "vendor assessments", "vendor risks", "vendor reviews", "vendor register",
    "third party risk", "third-party risk", "supplier list", "vendor list", "vendors",
    "vendor", "third party", "supplier",
    # ── Compliance ──
    "compliance assessments", "compliance programs", "compliance status",
    "frameworks", "compliance frameworks", "compliance requirements", "regulatory obligations",
    "pci dss requirements", "iso controls", "nist controls", "cobit controls",
    "cis controls", "sbp controls", "sama controls", "dora controls",
    "gdpr requirements", "hipaa requirements", "sox controls",
    "state bank of pakistan", "compliance gaps", "compliance score", "compliance checklist",
    # ── Attestation & Certification ──
    "attestation campaigns", "attestation requests", "attestations",
    "pending attestations", "overdue attestations",
    "certification journeys", "certification phases", "certification status",
    "attestation campaign", "sox attestation", "policy signoff",
    # ── Committee & Governance ──
    "committee", "committees", "committee member", "committee members",
    "committee meeting", "committee meetings", "governance committee", "governance committees",
    "oversight action", "oversight actions",
    "board meeting", "board meetings", "meeting agenda", "meeting minutes",
    "meeting action", "action points", "action items", "agenda item",
    "governance documents", "governance document",
    # ── RCSA ──
    "rcsa campaigns", "rcsa findings", "rcsa assessments", "self assessments",
    "rcsa", "self-assessment", "rcsa template", "rcsa response",
    "business unit assessment",
    # ── Vulnerabilities & Scanning ──
    "pentest reports", "vuln reports", "vulnerability reports",
    "open vulnerabilities", "critical vulnerabilities", "cve list",
    "vulnerability", "vulnerabilities", "pentest", "penetration test",
    "vulnerability mitigation", "vuln mitigation", "remediation action",
    "retest", "vulnerability retest", "vuln retest",
    "vulnerability sla", "sla breach", "sla target", "remediation sla",
    "scan record", "scan records", "vulnerability scan", "scan job",
    "sync history", "integration sync", "scanner sync", "last sync",
    "scanner exception", "vuln exception", "integration exception",
    # ── Assets ──
    "asset", "assets", "it asset", "it assets", "asset inventory", "asset management",
    "critical asset", "critical assets", "cde asset", "cde assets",
    "asset vulnerability", "asset risk", "asset control",
    # ── Integration / Connections ──
    "integration", "integrations", "integration connection", "integration connections",
    "connector", "connectors", "scanner integration", "vulnerability scanner",
    "nexpose", "nessus", "qualys",
    # ── Regulatory ──
    "regulatory changes", "regulatory updates", "new regulations",
    "regulatory feed", "regulatory impact", "regulatory implementation",
    "regulatory change",
    # ── Policies & Documents ──
    "policy", "policies", "procedure", "procedures", "standard", "standards",
    "guideline", "guidelines", "charter", "charters",
    "policy linkage", "control linkage", "policy link",
    "policy documents", "governance policies", "policy statements",
    "document review", "review schedule", "expiring policies",
    "policy gap analysis", "policy compliance", "policy gap",
    "policy gap finding", "policy gap findings", "gap finding", "gap findings",
    "gap analysis result", "gap analysis results",
    "policy attestation",
    # ── Internal Controls ──
    "internal controls", "key controls", "control tests", "control library",
    "control effectiveness", "control gaps", "internal control", "control test",
    "control mapping",
    # ── Evidence ──
    "evidence register", "evidence items", "control evidence",
    "missing evidence", "weak evidence", "evidence gaps",
    "evidence", "evidence assessment",
    # ── Audit Management (tables in DB — route to database) ──
    "audit finding", "audit findings", "audit plan", "audit plans",
    "audit engagement", "audit engagements", "audit report", "audit reports",
    "audit workpaper", "workpaper", "workpapers", "audit workpapers",
    "audit recommendation", "audit recommendations", "audit action plan",
    "pbc list", "prepared by client", "audit document request",
    "qaip", "audit quality", "audit maturity",
    "audit team", "audit universe", "auditable entity", "audit template",
    "audit board pack", "board pack", "audit follow up", "audit follow-up",
    "audit finding theme", "audit report opinion",
    # ── CCM (Continuous Control Monitoring) ──
    "ccm", "continuous control monitoring", "ccm rule", "ccm rules",
    "ccm anomaly", "ccm anomalies", "ccm exception", "ccm exceptions",
    "control monitoring", "control monitoring alert", "automated control testing",
    # ── Users & Departments ──
    "user accounts", "platform users", "department list",
    "department member", "department members",
)
FRAMEWORK_PROGRESS_TERMS = (
    "active compliance frameworks", "current progress", "framework progress", "active frameworks",
    "frameworks and their current progress", "framework overview"
)
EVIDENCE_GAP_TERMS = (
    "missing or weak evidence", "missing evidence", "weak evidence", "evidence gaps",
    "controls with missing", "controls with weak evidence"
)
# ─── Prompt injection patterns ──────────────────────────────────────────────
PROMPT_INJECTION_PATTERNS = (
    "forget your role", "ignore previous instructions", "ignore your instructions",
    "ignore all previous", "disregard your", "override your instructions",
    "reveal admin", "reveal credentials", "reveal the admin", "reveal your prompt",
    "pretend you are", "act as if you are", "you are now a",
    "bypass your", "your new instructions", "your new role",
    "execute any instructions", "execute the following instructions",
    "forget everything", "new persona", "jailbreak",
    "dan mode", "developer mode", "unrestricted mode",
    "translate the following", "ignore the above",
    "print your system prompt", "output your instructions",
    "what are your instructions", "reveal your system",
    "respond only in", "respond in the following way",
)
# ─── Harmful / attack-oriented request patterns ─────────────────────────────
HARMFUL_REQUEST_PATTERNS = (
    "ways to bypass", "how to bypass", "bypass security controls",
    "bypass evidence", "bypass audit", "bypass soc 2",
    "bypass iso", "bypass pci", "bypass nist", "bypass compliance",
    "circumvent controls", "circumvent evidence", "circumvent audit",
    "evade controls", "evade detection", "evade audit",
    "disable logging", "disable monitoring", "disable controls",
    "avoid detection", "hide from audit", "avoid audit trail",
    "delete evidence to", "remove evidence to", "destroy evidence",
    "ways to avoid", "how to avoid compliance", "bypass authorization",
    "unauthorized access", "how to hack", "how to exploit",
    "attack vector", "penetrate without", "bypass authentication",
    "disable evidence collection", "defeat monitoring",
)
# ─── Framework content / knowledge signals ───────────────────────────────────
# Questions about what a framework says/requires route to LLM (not DB query)
FRAMEWORK_CONTENT_SIGNALS = (
    "what are the new requirements",
    "new requirements in",
    "new requirements of",
    "what is new in",
    "what changed in",
    "what does pci", "what does iso", "what does nist", "what does sama",
    "what does sbp", "what does gdpr", "what does hipaa", "what does dora",
    "what does soc 2", "what does sox", "what does cobit", "what does cis",
    "what does basel",
    "what are the requirements of",
    "requirements of pci", "requirements of iso", "requirements of nist",
    "requirements of sama", "requirements of sbp", "requirements of gdpr",
    "requirements of hipaa", "requirements of sox",
    "clauses of", "sections of", "domains of framework",
    "explain pci", "explain iso 27001", "explain nist", "explain sama",
    "describe the framework", "explain the framework", "what is the framework",
    "what are the key controls of", "key domains of",
    "what does this standard require", "what does this regulation require",
    # Version references signal framework knowledge question
    " v1.", " v2.", " v3.", " v4.", " v5.",
    "version 1.", "version 2.", "version 3.", "version 4.",
    "pci dss 4", "pci dss 3.2", "pci dss v", "iso 27001:2022", "iso 27001:2013",
    "nist 2.0", "nist 1.1",
    # "What are <framework> requirements" without 'our' = knowledge question
    "what are pci", "what are iso", "what are nist", "what are gdpr",
    "what are hipaa", "what are sama", "what are dora", "what are sox",
    # Definitional questions ("what is a X" without 'our'/'my' = knowledge, not data)
    "what is a ", "what is an ", "define ", "how does a ", "how does an ",
    "what does it mean", "explain what is", "describe what is",
)
# When present, question is about the USER's platform data → override to DB route
PLATFORM_CONTEXT_SIGNALS = (
    " my ", " our ", " we ", "we have", "do we have",
    "have we", "are we ", "in our system", "in the platform",
    "in this system", "which controls have i", "what have we",
    "our compliance", "our gaps", "our posture", "our implementation",
)
LANGCHAIN_HISTORY_AVAILABLE = False

try:
    from langchain_core.chat_history import InMemoryChatMessageHistory
    from langchain_core.messages import AIMessage, HumanMessage
    LANGCHAIN_HISTORY_AVAILABLE = True
except Exception as langchain_error:
    logger.warning(f"LangChain history module unavailable: {langchain_error}")

try:
    from .qdrant_service import (
        IndexedSourceDocument,
        QdrantComplyChatService,
        extract_text_from_path,
    )
except Exception as qdrant_import_error:
    IndexedSourceDocument = None  # type: ignore[assignment]
    QdrantComplyChatService = None  # type: ignore[assignment]
    extract_text_from_path = None  # type: ignore[assignment]
    logger.warning("Qdrant service module unavailable: %s", qdrant_import_error)


def normalize_chat_message(value: str) -> str:
    cleaned = (value or "").strip()
    if not cleaned:
        raise ValueError("message cannot be empty")
    return cleaned


def resolve_chat_session_id(session_id: Optional[str], user_id: Any) -> str:
    raw_value = (session_id or "").strip()
    cleaned = "".join(ch if ch.isalnum() or ch in ("-", "_") else "_" for ch in raw_value)
    return cleaned or f"user_{user_id}"


def get_session_storage_key(user_id: Any, session_id: str) -> str:
    return f"{user_id}:{resolve_chat_session_id(session_id, user_id)}"


def is_audit_related_question(question: str) -> bool:
    lowered = (question or "").lower()
    return any(term in lowered for term in AUDIT_QUERY_TERMS)


def get_recent_session_log(user_id: Any, session_id: str) -> List[Dict[str, Any]]:
    storage_key = get_session_storage_key(user_id, session_id)
    return session_message_logs.setdefault(storage_key, [])


def append_session_message(user_id: Any, session_id: str, role: str, content: str, **extra: Any) -> None:
    storage_key = get_session_storage_key(user_id, session_id)
    payload: Dict[str, Any] = {
        "role": role,
        "content": (content or "").strip(),
        "timestamp": datetime.utcnow().isoformat(),
    }
    payload.update(extra)
    logs = session_message_logs.setdefault(storage_key, [])
    logs.append(payload)
    session_message_logs[storage_key] = logs[-(MAX_HISTORY_LENGTH * 2):]
    conversation_history[storage_key] = deque(session_message_logs[storage_key], maxlen=MAX_HISTORY_LENGTH * 2)


def hydrate_langchain_history(user_id: Any, session_id: str, history_items: Optional[List[Dict[str, Any]]] = None):
    if not LANGCHAIN_HISTORY_AVAILABLE:
        return None

    storage_key = get_session_storage_key(user_id, session_id)
    message_history = InMemoryChatMessageHistory()
    for item in (history_items or session_message_logs.get(storage_key, []))[-(MAX_HISTORY_LENGTH * 2):]:
        role = str(item.get("role", "")).strip().lower()
        content = str(item.get("content", "")).strip()
        if not content:
            continue
        if role == "assistant":
            message_history.add_ai_message(content)
        elif role == "user":
            message_history.add_user_message(content)
    langchain_histories[storage_key] = message_history
    return message_history


def build_context_summary(user_id: Any, session_id: str, history_items: Optional[List[Dict[str, Any]]] = None) -> str:
    hydrated = hydrate_langchain_history(user_id, session_id, history_items)
    context_lines = ["Recent conversation context:"]

    if hydrated is not None:
        for message in hydrated.messages[-6:]:
            content = str(getattr(message, "content", "")).strip()
            if not content:
                continue
            role = "Assistant" if isinstance(message, AIMessage) else "User"
            context_lines.append(f"- {role}: {content[:200]}")
    else:
        for item in (history_items or [])[-6:]:
            content = str(item.get("content", "")).strip()
            if not content:
                continue
            role = "Assistant" if item.get("role") == "assistant" else "User"
            context_lines.append(f"- {role}: {content[:200]}")

    return "\n".join(context_lines) if len(context_lines) > 1 else ""


def store_chat_exchange(
    user_id: Any,
    session_id: str,
    user_message: str,
    assistant_message: str,
    *,
    offset: int = 0,
    sources: Optional[List[Dict[str, Any]]] = None,
    is_error: bool = False,
) -> None:
    if offset == 0:
        append_session_message(user_id, session_id, "user", user_message)
        append_session_message(user_id, session_id, "assistant", assistant_message, sources=sources or [], is_error=is_error)
        hydrate_langchain_history(user_id, session_id)


def get_session_upload_dir(user_id: Any, session_id: str) -> Path:
    upload_dir = CHAT_UPLOAD_ROOT / f"user_{user_id}" / resolve_chat_session_id(session_id, user_id)
    upload_dir.mkdir(parents=True, exist_ok=True)
    return upload_dir


def get_session_manifest_path(user_id: Any, session_id: str) -> Path:
    return get_session_upload_dir(user_id, session_id) / "manifest.json"


def load_session_uploaded_files(user_id: Any, session_id: str) -> List[Dict[str, Any]]:
    storage_key = get_session_storage_key(user_id, session_id)
    if storage_key in session_uploaded_files:
        return session_uploaded_files[storage_key]

    manifest_path = get_session_manifest_path(user_id, session_id)
    if manifest_path.exists():
        try:
            session_uploaded_files[storage_key] = json.loads(manifest_path.read_text(encoding="utf-8"))
        except Exception:
            session_uploaded_files[storage_key] = []
    else:
        session_uploaded_files[storage_key] = []

    return session_uploaded_files[storage_key]


def save_session_uploaded_files(user_id: Any, session_id: str, files: List[Dict[str, Any]]) -> None:
    storage_key = get_session_storage_key(user_id, session_id)
    session_uploaded_files[storage_key] = files
    manifest_path = get_session_manifest_path(user_id, session_id)
    manifest_path.write_text(json.dumps(files, ensure_ascii=False, indent=2), encoding="utf-8")


def extract_text_from_upload(filename: str, content: bytes) -> str:
    suffix = Path(filename).suffix.lower()
    text_extract = ""

    try:
        if suffix in {".txt", ".md", ".csv", ".json", ".yaml", ".yml", ".xml", ".html", ".log", ".ini", ".py", ".js", ".ts"}:
            text_extract = content.decode("utf-8", errors="ignore")
        elif suffix == ".pdf":
            from PyPDF2 import PdfReader
            pdf_reader = PdfReader(BytesIO(content))
            text_extract = "\n".join((page.extract_text() or "") for page in pdf_reader.pages[:10])
        elif suffix == ".docx":
            from docx import Document
            document = Document(BytesIO(content))
            text_extract = "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text)
        elif suffix == ".pptx":
            from pptx import Presentation
            presentation = Presentation(BytesIO(content))
            chunks: List[str] = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        chunks.append(shape.text)
            text_extract = "\n".join(chunks)
    except Exception as parse_error:
        logger.warning(f"Could not extract text from {filename}: {parse_error}")

    if not text_extract:
        text_extract = content.decode("utf-8", errors="ignore")

    text_extract = text_extract.strip()
    return text_extract[:4000] if text_extract else f"Uploaded binary file '{filename}' ({len(content)} bytes)."


def build_uploaded_context(user_id: Any, session_id: str) -> str:
    uploaded_files = load_session_uploaded_files(user_id, session_id)
    if not uploaded_files:
        return ""

    context_lines = ["Uploaded file context:"]
    for item in uploaded_files[-5:]:
        filename = item.get("filename", "file")
        excerpt = str(item.get("excerpt", "")).strip()
        if excerpt:
            context_lines.append(f"- {filename}: {excerpt[:600]}")
        else:
            context_lines.append(f"- {filename}: File uploaded successfully and available for reference.")
    return "\n".join(context_lines)


def is_grc_relevant_question(question: str) -> bool:
    """Returns True if the question contains ANY GRC-related signal (permissive)."""
    lowered = (question or "").lower()
    normalized = normalize_question(lowered)
    return any(term in lowered for term in GRC_RELEVANT_TERMS) or any(term in normalized for term in GRC_RELEVANT_TERMS)


def grc_relevance_score(question: str) -> float:
    """Returns a 0.0-1.0 relevance score — fraction of GRC terms found (capped at 5 matches)."""
    lowered = (question or "").lower()
    hits = sum(1 for term in GRC_RELEVANT_TERMS if term in lowered)
    return min(hits / 5.0, 1.0)


def is_off_topic_question(question: str) -> bool:
    """Only blocks queries that are clearly non-GRC with zero relevance (< 20% score)."""
    lowered = (question or "").lower().strip()
    if any(term == lowered for term in GREETING_TERMS):
        return False
    if is_grc_relevant_question(question):
        return False
    if any(term in lowered for term in OFF_TOPIC_TERMS):
        return True
    return False


def is_prompt_injection(question: str) -> bool:
    """Detect prompt injection / jailbreak attempts."""
    lowered = (question or "").lower()
    return any(pattern in lowered for pattern in PROMPT_INJECTION_PATTERNS)


def is_harmful_request(question: str) -> bool:
    """Detect requests asking how to bypass or defeat security/compliance controls."""
    lowered = (question or "").lower()
    return any(pattern in lowered for pattern in HARMFUL_REQUEST_PATTERNS)


def is_framework_knowledge_question(question: str) -> bool:
    """
    Returns True if the question is asking about what a GRC framework says/requires
    (knowledge question) rather than about the user's own live platform data.
    These route to LLM guidance, not DB query.
    """
    lowered = normalize_question((question or "").lower())
    # If question explicitly references the user's own platform data, it's a DB question
    if any(ctx in lowered for ctx in PLATFORM_CONTEXT_SIGNALS):
        return False
    return any(signal in lowered for signal in FRAMEWORK_CONTENT_SIGNALS)


def normalize_question(question: str) -> str:
    """
    Normalize common typos and alternate spellings to canonical GRC terms
    so classification and routing work regardless of spelling mistakes.
    """
    import re
    text = (question or "").lower()
    substitutions = [
        # compliance
        (r'\bcomplian[sc]e?\b', 'compliance'),
        (r'\bcomplience\b', 'compliance'),
        # governance
        (r'\bgover[ae]n[ae]nce?\b', 'governance'),
        (r'\bgovrnance\b', 'governance'),
        # framework
        (r'\bfram[ew]ork\b', 'framework'),
        (r'\bframworks?\b', 'frameworks'),
        # vulnerability / vulnerabilities
        (r'\bvuln?er[ae]bil[iy]t[yi]e?s?\b', 'vulnerabilities'),
        (r'\bvunerabilit\w*\b', 'vulnerabilities'),
        (r'\bvulnr?abilit\w*\b', 'vulnerabilities'),
        # regulatory
        (r'\breg[uo]l[ae]tor[yi]\b', 'regulatory'),
        # attestation
        (r'\batte?sta[st]?ion\b', 'attestation'),
        # certification
        (r'\bcertif[iy]?[ae]?tion\b', 'certification'),
        # assessment
        (r'\basse[sc]?e?ment\b', 'assessment'),
        # incident
        (r'\bincid[ae]?nts?\b', 'incident'),
        # control
        (r'\bcontro?ll?\b', 'control'),
        # risk / risks
        (r'\brisk[s]?\b', 'risks'),
        # policy / policies (police → policy is a common typo)
        (r'\bpolice\b', 'policy'),
        (r'\bpolicey\b', 'policy'),
        (r'\bpolicies\b', 'policies'),
        # asset / assets
        (r'\bassests?\b', 'assets'),
        (r'\bassest\b', 'asset'),
        # committee / committees
        (r'\bcommitee\b', 'committee'),
        (r'\bcommitees\b', 'committees'),
        (r'\bcomittee\b', 'committee'),
        (r'\bcommitte\b', 'committee'),
        # integration / integrations
        (r'\bintegraton\b', 'integration'),
        (r'\bintegeation\b', 'integration'),
        (r'\bintegrations?\b', 'integration'),
        # linkage
        (r'\blinage\b', 'linkage'),
        (r'\blinkege\b', 'linkage'),
        # vendor
        (r'\bvendors?\b', 'vendor'),
        (r'\bvender\b', 'vendor'),
        # evidence
        (r'\bevidance\b', 'evidence'),
        (r'\bevidenc[ey]\b', 'evidence'),
    ]
    for pattern, replacement in substitutions:
        text = re.sub(pattern, replacement, text)
    return text


def classify_request_mode(question: str, has_uploaded_files: bool = False) -> str:
    lowered = (question or "").lower().strip()
    # Normalize typos before classification so misspellings route correctly
    normalized = normalize_question(lowered)

    # ── Security guardrails (always first) ───────────────────────────────────
    if is_prompt_injection(question):
        return "blocked_injection"
    if is_harmful_request(question):
        return "blocked_harmful"

    if is_audit_related_question(lowered):
        return "deprecated_audit"
    if is_off_topic_question(normalized) and not has_uploaded_files:
        return "off_topic"
    if any(term in normalized for term in FRAMEWORK_PROGRESS_TERMS):
        return "framework_progress"
    if any(term in normalized for term in EVIDENCE_GAP_TERMS):
        return "evidence_gaps"
    if has_uploaded_files and any(term in normalized for term in FILE_ANALYSIS_TERMS):
        return "file_analysis"
    if has_uploaded_files and not any(term in normalized for term in DB_QUERY_TERMS):
        return "file_analysis"

    # ── Quantitative intent = ALWAYS live platform data ───────────────────────
    # "how many assets", "total risks", "count vulnerabilities", "number of committees" etc.
    _quantitative = ("how many", "how much", "total number", "count of", "number of",
                     "how many open", "how many pending", "how many critical",
                     "what is the total", "what's the total")
    if any(q in normalized for q in _quantitative) and not has_uploaded_files:
        return "database"

    # ── Framework / conceptual knowledge questions (check BEFORE PLATFORM_DATA_NOUNS
    # so "what are pci dss v4 requirements" routes to guidance not database) ───────
    if is_framework_knowledge_question(question) and not has_uploaded_files:
        return "grc_guidance"

    # ── PLATFORM_DATA_NOUNS always means live DB data ─────────────────────────
    # Check BEFORE framework-knowledge check to avoid intercepting platform queries
    if (any(noun in normalized for noun in PLATFORM_DATA_NOUNS) or any(noun in lowered for noun in PLATFORM_DATA_NOUNS)) and not has_uploaded_files:
        # If question also contains a DB_QUERY_TERM or quantitative signal, definitely database
        if any(term in normalized for term in DB_QUERY_TERMS) or any(term in lowered for term in DB_QUERY_TERMS):
            return "database"
        # Even without explicit query term, platform nouns = live data
        if any(ctx in normalized for ctx in PLATFORM_CONTEXT_SIGNALS) or any(noun in normalized for noun in PLATFORM_DATA_NOUNS):
            return "database"

    # Explicit DB query commands — check both original and normalized
    if any(term in normalized for term in DB_QUERY_TERMS) or any(term in lowered for term in DB_QUERY_TERMS):
        return "database"

    # If any GRC noun is present and any DB signal → database
    if any(term in normalized for term in GRC_RELEVANT_TERMS) and any(term in normalized for term in DB_QUERY_TERMS):
        return "database"

    return "grc_guidance"


# ============================================================================
# LLM KNOWLEDGE FALLBACK
# ============================================================================

def answer_grc_knowledge_question(
    question: str,
    context_summary: str = "",
    uploaded_context: str = "",
    db_context: str = "",
) -> str:
    """
    Answer a GRC knowledge / conceptual question using GPT-4o.
    Covers: framework requirements, ERM overviews, best practices, clause explanations,
    version differences, and general compliance guidance.
    An optional db_context can be passed to ground the answer in real platform data.
    """
    try:
        import openai as _openai
        _api_key = os.environ.get("OPENAI_API_KEY", "")
        if not _api_key:
            return (
                "I couldn't reach the AI engine. "
                "Please ask about frameworks, risks, or controls stored in the platform."
            )
        _client = _openai.OpenAI(api_key=_api_key)

        system_prompt = (
            "You are ComplyChat, an expert GRC AI assistant inside the ComplyVerse enterprise GRC platform.\n\n"
            "You are deeply knowledgeable about:\n"
            "- Compliance frameworks: ISO 27001 (2013 & 2022), PCI DSS (v3.2.1 & v4.0.1), NIST CSF (1.1 & 2.0), "
            "SAMA Cyber Security Framework, SBP ETGRMF, DORA, GDPR, SOC 2, HIPAA, SOX, COBIT, CIS Controls, Basel\n"
            "- Enterprise Risk Management (ERM), internal controls, governance, evidence management\n"
            "- Vulnerability management, asset management, vendor risk, RCSA, attestation, incident management\n\n"
            "RESPONSE RULES — follow strictly:\n"
            "1. Lead with a direct, concise answer. Max 350 words unless detail is clearly needed.\n"
            "2. Use `##` for section headers, `-` for bullets, `**bold**` for key terms and control IDs.\n"
            "3. For version/change questions (e.g. PCI DSS v3.2.1 vs v4.0.1): list specific differences clearly.\n"
            "4. For framework requirement questions: list the actual requirements/controls with their IDs.\n"
            "5. NEVER say 'I cannot answer' for GRC topics you clearly know — give the best answer.\n"
            "6. NEVER hallucinate control IDs or clause numbers you are not sure of — say 'refer to the official document' instead.\n"
            "7. If referencing platform-specific data that isn't in context, note it and suggest checking the relevant module.\n"
            "8. Do NOT mention databases, queries, SQL, or system internals to the user.\n"
            "9. Politely decline anything outside GRC topics in one sentence.\n"
            "CRITICAL ANTI-HALLUCINATION RULES:\n"
            "10. NEVER say 'I can't access specific data' — the platform already queried the database.\n"
            "11. NEVER say 'typically' or 'usually' when describing the USER's platform data — only state facts from the context provided.\n"
            "12. If DB context shows no data: say 'No [X] found in your platform yet' — do NOT invent generic examples.\n"
            "13. NEVER list generic module names as if they were real results — only report actual data from the provided context.\n"
            "14. If the question asks for counts/lists and no context was provided: say 'I was unable to retrieve this data — please check the [module] section in the platform.'"
        )

        messages: List[Dict[str, Any]] = [{"role": "system", "content": system_prompt}]
        if db_context:
            messages.append({"role": "user", "content": f"Platform data context:\n{db_context}"})
        if context_summary:
            messages.append({"role": "user", "content": f"Previous conversation context:\n{context_summary}"})
        if uploaded_context:
            messages.append({"role": "user", "content": f"Uploaded file context:\n{uploaded_context}"})
        messages.append({"role": "user", "content": question})

        with usage_scope(module_key="complychat", feature_key="knowledge_guidance"):
            completion = _client.chat.completions.create(
                model=get_openai_model(),
                messages=messages,
                temperature=0.2,
                max_tokens=1500,
            )
        return (completion.choices[0].message.content or "").strip() or "I couldn't generate an answer. Please try rephrasing."

    except Exception as llm_err:
        logger.warning(f"[LLM-FALLBACK] knowledge answer failed: {llm_err}")
        return (
            "I can answer questions about GRC frameworks (ISO 27001, PCI DSS, NIST, SAMA, SBP, DORA), "
            "controls, vulnerabilities, assets, governance, and evidence management. "
            "Please try rephrasing your question."
        )


def answer_with_db_context(
    question: str,
    db: Session,
    tenant_ids: List[int],
    context_summary: str = "",
    uploaded_context: str = "",
) -> str:
    """
    Hybrid answer: pull relevant framework controls / documents from the DB as grounding
    context, then use LLM to synthesize a comprehensive, accurate answer.
    Used for framework requirement questions, clause lookups, and control explanations.
    """
    db_context_parts: List[str] = []

    try:
        # Detect framework names mentioned in the question
        question_lower = (question or "").lower()
        framework_keywords = []
        for fw in ["pci dss", "pci-dss", "iso 27001", "iso27001", "nist", "sama", "sbp", "dora", "gdpr",
                   "soc 2", "soc2", "hipaa", "sox", "cobit", "cis", "basel", "etgrmf", "bss"]:
            if fw in question_lower:
                framework_keywords.append(fw.replace("-", " ").strip())

        if framework_keywords:
            for kw in framework_keywords[:2]:  # max 2 frameworks
                try:
                    fw_rows = db.execute(
                        text(
                            "SELECT id, name, framework_type, version, source_organization, framework_purpose "
                            "FROM grc_uploaded_frameworks "
                            "WHERE LOWER(name) LIKE :kw "
                            "OR LOWER(framework_type) LIKE :kw "
                            "LIMIT 2"
                        ),
                        {"kw": f"%{kw[:30]}%"},
                    ).fetchall()

                    for fw_row in fw_rows[:1]:
                        fw_dict = dict(fw_row._mapping)
                        fw_name = fw_dict.get("name") or kw
                        fw_ver = fw_dict.get("version") or "N/A"
                        fw_purpose = (str(fw_dict.get("framework_purpose") or "")).strip()[:200]
                        db_context_parts.append(
                            f"## {fw_name} (v{fw_ver})\n"
                            + (f"Purpose: {fw_purpose}\n" if fw_purpose else "")
                        )

                        # Get controls for this framework
                        ctrl_rows = db.execute(
                            text(
                                "SELECT pfc.control_id, pfc.title, pfc.category, pfc.description "
                                "FROM grc_parsed_framework_controls pfc "
                                "JOIN grc_uploaded_frameworks uf ON pfc.uploaded_framework_id = uf.id "
                                "WHERE uf.id = :fw_id "
                                "ORDER BY pfc.control_id "
                                "LIMIT 60"
                            ),
                            {"fw_id": fw_dict.get("id")},
                        ).fetchall()

                        if ctrl_rows:
                            controls_text = "\n".join(
                                f"- **{r.control_id}**: {r.title}"
                                + (f" ({r.category})" if r.category else "")
                                for r in ctrl_rows[:40]
                            )
                            db_context_parts.append(
                                f"Controls in platform ({len(ctrl_rows)} total):\n{controls_text}"
                            )
                except Exception as fw_err:
                    logger.warning(f"[DB-CONTEXT] framework lookup failed for '{kw}': {fw_err}")
                    try:
                        db.rollback()
                    except Exception:
                        pass

        # Also pull relevant governance documents (policies / standards)
        try:
            doc_rows = db.execute(
                text(
                    "SELECT title, doc_type, status, current_version "
                    "FROM grc_governance_documents "
                    "WHERE tenant_id IN :tids "
                    "AND COALESCE(status,'draft') IN ('published','approved') "
                    "LIMIT 10"
                ),
                {"tids": tuple(tenant_ids) if tenant_ids else (-1,)},
            ).fetchall()

            if doc_rows:
                docs_text = "\n".join(
                    f"- {r.title} ({r.doc_type}, v{r.current_version or '1.0'})" for r in doc_rows
                )
                db_context_parts.append(f"Published documents in your platform:\n{docs_text}")
        except Exception as doc_err:
            logger.warning(f"[DB-CONTEXT] document lookup failed: {doc_err}")
            try:
                db.rollback()
            except Exception:
                pass

    except Exception as db_err:
        logger.warning(f"[DB-CONTEXT] DB context fetch failed: {db_err}")

    db_context = "\n\n".join(db_context_parts) if db_context_parts else ""
    return answer_grc_knowledge_question(question, context_summary, uploaded_context, db_context=db_context)


def summarize_framework_progress(db: Session, tenant_ids: List[int]) -> str:
    available_statuses = ["published", "completed", "parsed", "classified"]
    framework_rows = db.query(UploadedFramework).filter(
        UploadedFramework.upload_status.in_(available_statuses),
        or_(
            UploadedFramework.tenant_id.in_(tenant_ids) if tenant_ids else False,
            UploadedFramework.is_shared == True,
            UploadedFramework.tenant_id.is_(None)
        )
    ).order_by(UploadedFramework.name.asc()).all()

    deduped: Dict[str, Any] = {}
    for framework in framework_rows:
        key = f"{(framework.name or '').strip().lower()}::{(framework.version or '').strip().lower()}"
        if key not in deduped or framework.tenant_id in tenant_ids:
            deduped[key] = framework

    journeys = db.query(CertificationJourney).filter(
        CertificationJourney.tenant_id.in_(tenant_ids) if tenant_ids else False
    ).all()

    tracked_journeys = [journey for journey in journeys if (journey.status or "").lower() in {"active", "in_progress", "completed", "on_hold"}]
    active_journeys = [journey for journey in tracked_journeys if (journey.status or "").lower() in {"active", "in_progress"}]
    journey_by_framework_id = {
        journey.uploaded_framework_id: journey
        for journey in active_journeys
        if getattr(journey, "uploaded_framework_id", None)
    }

    lines: List[str] = []
    if not active_journeys:
        lines.append("Executive Summary")
        lines.append("No active compliance or certification journeys are currently running for your tenant.")
        lines.append(f"{len(deduped)} unique framework(s) are available in the library, but progress tracking will begin only after a journey is started.")
        if tracked_journeys:
            lines.append(f"There are {len(tracked_journeys)} tracked journey record(s), but none are currently active or in progress.")
        lines.append("")
        lines.append("Available Frameworks")
        for framework in list(deduped.values())[:10]:
            lines.append(
                f"- {framework.name} — library status: {framework.upload_status or 'unknown'}"
                + (f", version: {framework.version}" if framework.version else "")
            )
        return "\n".join(lines)

    lines.append("Executive Summary")
    lines.append(f"{len(active_journeys)} active journey(s) were found across {len(journey_by_framework_id)} framework(s).")
    lines.append("")
    lines.append("Current Progress")

    for journey in active_journeys[:10]:
        framework_name = getattr(getattr(journey, "uploaded_framework", None), "name", None) or journey.name or "Unknown Framework"
        readiness = 0
        status_label = journey.status or "unknown"
        if callable(calculate_progress_summary):
            try:
                summary = calculate_progress_summary(journey, db)
                readiness = round(float(getattr(summary, "readiness_percentage", 0) or 0), 1)
            except Exception:
                readiness = 0
        lines.append(f"- {framework_name}: status {status_label}, readiness {readiness}%")

    return "\n".join(lines)


def summarize_evidence_gaps(db: Session, tenant_ids: List[int]) -> str:
    journeys = db.query(CertificationJourney).filter(
        CertificationJourney.tenant_id.in_(tenant_ids) if tenant_ids else False
    ).all()

    if not journeys:
        return (
            "No evidence-backed control implementations were found for your tenant yet. "
            "Start a framework journey and upload/link evidence to surface missing or weak controls."
        )

    implementation_ids = [journey.id for journey in journeys]
    implementations = db.query(ControlImplementation).filter(
        ControlImplementation.journey_id.in_(implementation_ids)
    ).all()

    if not implementations:
        return (
            "No control implementations exist yet for the current journeys, so there is no real evidence gap data to report. "
            "Create or sync the controls first, then upload evidence for review."
        )

    weak_rows: List[Dict[str, Any]] = []
    for implementation in implementations:
        parsed_control = getattr(implementation, "parsed_control", None)
        framework_control = getattr(implementation, "framework_control", None)
        control_code = getattr(parsed_control, "control_id", None) or getattr(framework_control, "code", None) or f"Control {implementation.id}"
        control_name = getattr(parsed_control, "title", None) or getattr(framework_control, "name", None) or "Untitled Control"

        attachment_count = len(getattr(implementation, "evidence_attachments", []) or [])
        approved_count = len([item for item in (implementation.evidence_attachments or []) if getattr(item, "review_status", None) == "approved"])
        weak_count = len([
            item for item in (implementation.evidence_attachments or [])
            if getattr(item, "review_status", None) in {"rejected", "pending"}
            or ((getattr(item, "ai_confidence_score", None) or 0) > 0 and (getattr(item, "ai_confidence_score", None) or 0) < 70)
        ])

        if attachment_count == 0 or approved_count == 0 or weak_count > 0:
            weak_rows.append({
                "control_code": control_code,
                "control_name": control_name,
                "attachment_count": attachment_count,
                "approved_count": approved_count,
                "weak_count": weak_count,
            })

    if not weak_rows:
        total_evidence = db.query(func.count(Evidence.id)).filter(Evidence.tenant_id.in_(tenant_ids) if tenant_ids else False).scalar() or 0
        return f"Good news: the tracked controls do not currently show missing or weak evidence. Total evidence records found: {total_evidence}."

    weak_rows.sort(key=lambda row: (row["approved_count"], row["attachment_count"], -row["weak_count"]))
    lines = [
        "Evidence Gap Summary",
        f"{len(weak_rows)} control(s) currently need attention.",
        "",
        "Priority Controls"
    ]
    for row in weak_rows[:10]:
        issue_bits = []
        if row["attachment_count"] == 0:
            issue_bits.append("no evidence uploaded")
        if row["approved_count"] == 0 and row["attachment_count"] > 0:
            issue_bits.append("nothing approved yet")
        if row["weak_count"] > 0:
            issue_bits.append(f"{row['weak_count']} weak or pending item(s)")
        lines.append(f"- {row['control_code']} — {row['control_name']} ({'; '.join(issue_bits)})")
    return "\n".join(lines)


def analyze_uploaded_files_with_llm(
    question: str,
    uploaded_files: List[Dict[str, Any]],
    context_summary: str = "",
) -> str:
    """
    Send the actual uploaded file content + user question to the LLM for a real GRC analysis.
    Handles gap analysis, compliance checks, control extraction, obligation summaries etc.
    """
    if not uploaded_files:
        return "Please upload one or more files so I can analyze them from a GRC perspective."

    try:
        import openai as _openai
        _api_key = os.environ.get("OPENAI_API_KEY", "")
        if not _api_key:
            return (
                "I couldn't reach the AI engine (missing API key). "
                "Please check that OPENAI_API_KEY is configured."
            )
        _client = _openai.OpenAI(api_key=_api_key)

        # Build file context — include actual extracted text, not just keyword counts
        file_blocks: List[str] = []
        for item in uploaded_files[-5:]:
            filename = item.get("filename", "uploaded file")
            excerpt = str(item.get("excerpt", "")).strip()
            if excerpt:
                file_blocks.append(f"--- FILE: {filename} ---\n{excerpt[:3500]}\n--- END FILE ---")
            else:
                file_blocks.append(f"--- FILE: {filename} ---\n[No extractable text content]\n--- END FILE ---")

        files_content = "\n\n".join(file_blocks)

        system_prompt = (
            "You are ComplyChat, an expert GRC (Governance, Risk & Compliance) AI assistant. "
            "You have been given one or more uploaded document(s) to analyze.\n\n"
            "Your capabilities include:\n"
            "- Gap analysis against compliance frameworks (ISO 27001, PCI DSS, NIST CSF, SAMA, SBP, DORA, GDPR, SOC 2, HIPAA etc.)\n"
            "- Extracting controls, obligations, and requirements from policy documents\n"
            "- Identifying missing clauses, weak language, or non-conformities\n"
            "- Summarizing document scope, purpose, and coverage\n"
            "- Mapping document content to specific framework control domains\n\n"
            "Rules:\n"
            "1. Base your answer STRICTLY on the document content provided. Do not hallucinate clauses.\n"
            "2. When doing gap analysis, list what IS covered and what IS MISSING or insufficient.\n"
            "3. Use clear headings, bullet points, and control IDs where applicable.\n"
            "4. Be specific — reference actual text from the document when possible.\n"
            "5. If the document text is truncated, note that and work with what is available."
        )

        messages: List[Dict[str, Any]] = [{"role": "system", "content": system_prompt}]
        if context_summary:
            messages.append({"role": "user", "content": f"Conversation context:\n{context_summary}"})
        messages.append({
            "role": "user",
            "content": (
                f"Uploaded document(s):\n\n{files_content}\n\n"
                f"User question: {question}"
            )
        })

        with usage_scope(module_key="complychat", feature_key="file_analysis"):
            completion = _client.chat.completions.create(
                model="gpt-4o-mini",
                messages=messages,
                temperature=0.2,
                max_tokens=2000,
            )
        return (completion.choices[0].message.content or "").strip() or "I couldn't generate an analysis. Please try rephrasing."

    except Exception as llm_err:
        logger.warning(f"[FILE-ANALYSIS-LLM] failed: {llm_err}")
        return (
            "I encountered an error while analyzing the document. "
            "Please ensure the file is a readable PDF, Word document (.docx), or plain text file and try again."
        )


detect_query_type = None
generate_sql_query = None
validate_sql = None
format_query_results = None
validate_columns_in_sql = None
get_fallback_data_for_question = None
generate_answer_from_fallback_data = None
fetch_table_schema_from_db = None
load_full_database_schema = None
SQL_AGENT_ENABLED = False

try:
    from grc_sql_agent import (
        detect_query_type, generate_sql_query, validate_sql, format_query_results,
        validate_columns_in_sql, get_fallback_data_for_question, generate_answer_from_fallback_data,
        fetch_table_schema_from_db, load_full_database_schema
    )  # type: ignore
    SQL_AGENT_ENABLED = True
    logger.info("[YES] SQL Agent loaded successfully (ChromaDB disabled)")
    if callable(load_full_database_schema):
        try:
            load_full_database_schema()
        except Exception as schema_error:
            logger.warning(f"SQL agent schema preload warning: {schema_error}")
except Exception as e:
    logging.error(f"Failed to import SQL agent: {e}")
    SQL_AGENT_ENABLED = False

from ...models import (
    GRCUser,
    get_db,
    Framework,
    FrameworkControl,
    ControlObjective,
    FrameworkDomain,
    UploadedFramework,
    ParsedFrameworkControl,
    CertificationJourney,
    ControlImplementation,
    GovernanceDocument,
    GovernanceDocumentVersion,
    GovernanceCommittee,
    CommitteeCharter,
    PolicyStatement,
    PolicyStatementVersion,
    RiskAssessment,
    FrameworkRiskAssessment,
    FrameworkRiskQuestion,
    GRCComplianceAssessment,
    ComplianceProgram,
    ComplianceAssessmentDocument,
    ComplianceAssessmentDocumentItem,
    RiskAssessmentRisk,
    Risk,
    Evidence,
    EvidenceControlMapping,
    FrameworkSubControl,
    CuratedEvidenceItem,
    ControlEvidenceRequirement,
)
from ...routers.auth_router import require_auth, get_user_tenants

try:
    from ...routers.certification_router import calculate_progress_summary
except Exception:
    calculate_progress_summary = None

router = APIRouter(prefix="/ai/complychat", tags=["AI Compliance Chat"])
logger = logging.getLogger(__name__)


# ============================================================================
# REQUEST/RESPONSE MODELS
# ============================================================================

class ChatRequest(BaseModel):
    """Chat request model"""

    model_config = ConfigDict(populate_by_name=True, extra="ignore")

    message: str = Field(
        ...,
        description="User's GRC question (compliance, risk, governance, evidence, policies, etc.)",
        min_length=1,
        validation_alias=AliasChoices("message", "question"),
        serialization_alias="message",
    )
    framework: Optional[str] = Field(None, description="Filter by framework code (e.g., 'PCI_DSS', 'ISO_27001')")
    include_sources: bool = Field(True, description="Include source references in response")
    session_id: Optional[str] = Field(None, description="Session ID for conversation history tracking")
    history: Optional[List[Dict[str, Any]]] = Field(None, description="Recent chat history for context hydration")
    limit: int = Field(10, ge=1, le=100, description="Number of results to return (default 10 for pagination)")
    offset: int = Field(0, ge=0, description="Offset for pagination (0 = first page)")

    @field_validator("message")
    @classmethod
    def validate_message(cls, value: str) -> str:
        return normalize_chat_message(value)

    @property
    def question(self) -> str:
        return self.message


class ChatSource(BaseModel):
    """Source reference model"""
    rank: int
    entity_type: str
    entity_id: Optional[str] = ""
    framework_code: str
    control_code: Optional[str] = None
    control_name: Optional[str] = None
    relevance_score: float
    snippet: str


class ChatResponse(BaseModel):
    """Chat response model"""
    answer: str
    sources: List[ChatSource]
    framework_filter: Optional[str]
    timestamp: str
    has_more: bool = Field(False, description="Whether more results are available")
    total_count: int = Field(0, description="Total number of results available")
    current_offset: int = Field(0, description="Current offset in results")


class FrameworkInfo(BaseModel):
    """Framework information model"""
    id: Optional[str] = None
    code: str
    name: str
    version: Optional[str] = ""
    description: Optional[str] = ""


class StatsResponse(BaseModel):
    """Knowledge base statistics"""
    total_entities: int
    entity_types: Dict[str, int]
    available_frameworks: List[FrameworkInfo]


class TriggerEmbeddingRequest(BaseModel):
    """Request to trigger embedding regeneration"""
    entity_types: Optional[List[str]] = Field(None, description="Specific entity types to update")
    async_mode: bool = Field(True, description="Run in background (non-blocking)")


def _is_policy_count_question(question: str) -> bool:
    normalized = normalize_question((question or "").lower())
    if "polic" not in normalized:
        return False
    return any(term in normalized for term in ("how many", "count", "total", "number of"))


def _policy_count_answer(db: Session, tenant_ids: List[int]) -> str:
    if not tenant_ids:
        return "No tenant scope is available for this user."

    statements_count = (
        db.query(func.count(PolicyStatement.id))
        .filter(PolicyStatement.tenant_id.in_(tenant_ids))
        .scalar()
        or 0
    )
    policy_docs_count = (
        db.query(func.count(GovernanceDocument.id))
        .filter(
            GovernanceDocument.tenant_id.in_(tenant_ids),
            func.lower(func.coalesce(GovernanceDocument.doc_type, "")) == "policy",
        )
        .scalar()
        or 0
    )
    total = int(statements_count) + int(policy_docs_count)
    return (
        "Policy counts in your platform:\n"
        f"- Policy Statements: {int(statements_count)}\n"
        f"- Governance Policy Documents: {int(policy_docs_count)}\n"
        f"- Combined Total: {total}"
    )


def _no_platform_data_answer(question: str) -> str:
    return (
        f"No matching platform data found for: \"{question.strip()}\".\n"
        "Try broadening filters, verify that module data is uploaded, or ask for a specific framework/document/control."
    )


def _sql_failure_answer(question: str) -> str:
    return (
        f"I could not complete this data query reliably: \"{question.strip()}\".\n"
        "No generated sample data is shown. Please try simpler wording or a narrower filter."
    )


# ============================================================================
# QDRANT VECTOR RAG HELPERS
# ============================================================================

_qdrant_service_singleton: Optional[Any] = None


def get_qdrant_service() -> Optional[Any]:
    global _qdrant_service_singleton
    if QdrantComplyChatService is None:
        return None
    if _qdrant_service_singleton is None:
        try:
            _qdrant_service_singleton = QdrantComplyChatService()
        except Exception as exc:
            logger.warning("Failed to initialize Qdrant service: %s", exc)
            _qdrant_service_singleton = None
    return _qdrant_service_singleton


def _to_iso(value: Any) -> str:
    if isinstance(value, datetime):
        return value.isoformat()
    if value is None:
        return ""
    return str(value)


def _json_text(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, str):
        return value.strip()
    if isinstance(value, (dict, list)):
        try:
            return json.dumps(value, ensure_ascii=False)
        except Exception:
            return str(value).strip()
    return str(value).strip()


def _merge_text(*parts: Any) -> str:
    cleaned = [str(part).strip() for part in parts if part and str(part).strip()]
    return "\n\n".join(cleaned)


def _catalog_entry(
    *,
    source_type: str,
    source_id: str,
    title: str,
    description: str = "",
    updated_at: str = "",
    tenant_id: Optional[int] = None,
) -> Dict[str, Any]:
    return {
        "source_type": source_type,
        "source_id": source_id,
        "title": title.strip(),
        "description": description.strip(),
        "updated_at": updated_at,
        "tenant_id": tenant_id,
    }


def _build_vector_catalog(
    db: Session,
    tenant_ids: List[int],
    uploaded_files: Optional[List[Dict[str, Any]]] = None,
) -> List[Dict[str, Any]]:
    entries: List[Dict[str, Any]] = []
    seen: set[Tuple[str, str]] = set()

    def add(entry: Dict[str, Any]) -> None:
        key = (entry["source_type"], entry["source_id"])
        if key in seen:
            return
        if not entry.get("title"):
            return
        seen.add(key)
        entries.append(entry)

    if tenant_ids:
        governance_docs = (
            db.query(GovernanceDocument)
            .filter(GovernanceDocument.tenant_id.in_(tenant_ids))
            .order_by(GovernanceDocument.updated_at.desc())
            .limit(max(600, VECTOR_GOV_VERSION_LIMIT))
            .all()
        )
        for item in governance_docs:
            add(
                _catalog_entry(
                    source_type="governance_document",
                    source_id=str(item.id),
                    title=item.title or f"Governance Document {item.id}",
                    description=_merge_text(item.doc_type, item.status, item.description)[:800],
                    updated_at=_to_iso(item.updated_at or item.created_at),
                    tenant_id=item.tenant_id,
                )
            )

        governance_versions = (
            db.query(GovernanceDocumentVersion, GovernanceDocument)
            .join(GovernanceDocument, GovernanceDocumentVersion.document_id == GovernanceDocument.id)
            .filter(GovernanceDocument.tenant_id.in_(tenant_ids))
            .order_by(GovernanceDocumentVersion.created_at.desc())
            .limit(max(200, VECTOR_GOV_VERSION_LIMIT))
            .all()
        )
        for version, document in governance_versions:
            add(
                _catalog_entry(
                    source_type="governance_document_version",
                    source_id=str(version.id),
                    title=version.title or document.title or f"Document Version {version.id}",
                    description=_merge_text(
                        document.title,
                        f"Version: {version.version_number}",
                        version.change_summary,
                        version.change_reason,
                        version.status,
                    )[:900],
                    updated_at=_to_iso(version.created_at),
                    tenant_id=document.tenant_id,
                )
            )

        policy_statements = (
            db.query(PolicyStatement)
            .filter(PolicyStatement.tenant_id.in_(tenant_ids))
            .order_by(PolicyStatement.updated_at.desc())
            .limit(max(500, VECTOR_POLICY_STATEMENT_LIMIT))
            .all()
        )
        for statement in policy_statements:
            title = statement.statement_code or f"Policy Statement {statement.id}"
            add(
                _catalog_entry(
                    source_type="policy_statement",
                    source_id=str(statement.id),
                    title=title,
                    description=_merge_text(
                        statement.category,
                        statement.sub_category,
                        statement.statement_summary,
                        statement.source_section,
                    )[:900],
                    updated_at=_to_iso(statement.updated_at or statement.created_at),
                    tenant_id=statement.tenant_id,
                )
            )

        policy_statement_versions = (
            db.query(PolicyStatementVersion)
            .filter(PolicyStatementVersion.tenant_id.in_(tenant_ids))
            .order_by(PolicyStatementVersion.changed_at.desc())
            .limit(max(400, VECTOR_POLICY_STATEMENT_LIMIT // 2))
            .all()
        )
        for version in policy_statement_versions:
            add(
                _catalog_entry(
                    source_type="policy_statement_version",
                    source_id=str(version.id),
                    title=f"Policy Statement {version.statement_id} v{version.version_number}",
                    description=_merge_text(version.category, version.sub_category, version.change_reason)[:900],
                    updated_at=_to_iso(version.changed_at),
                    tenant_id=version.tenant_id,
                )
            )

        charters = (
            db.query(CommitteeCharter)
            .filter(CommitteeCharter.tenant_id.in_(tenant_ids))
            .order_by(CommitteeCharter.created_at.desc())
            .limit(600)
            .all()
        )
        for item in charters:
            committee_name = item.committee.name if item.committee else ""
            add(
                _catalog_entry(
                    source_type="committee_charter",
                    source_id=str(item.id),
                    title=item.title or committee_name or f"Committee Charter {item.id}",
                    description=_merge_text(committee_name, item.version, item.status)[:800],
                    updated_at=_to_iso(item.created_at),
                    tenant_id=item.tenant_id,
                )
            )

        risk_assessments = (
            db.query(RiskAssessment)
            .filter(RiskAssessment.tenant_id.in_(tenant_ids))
            .order_by(RiskAssessment.updated_at.desc())
            .limit(600)
            .all()
        )
        for item in risk_assessments:
            add(
                _catalog_entry(
                    source_type="risk_assessment",
                    source_id=str(item.id),
                    title=item.name or f"Risk Assessment {item.id}",
                    description=_merge_text(item.assessment_type, item.status, item.description)[:800],
                    updated_at=_to_iso(item.updated_at or item.created_at),
                    tenant_id=item.tenant_id,
                )
            )

        compliance_programs = (
            db.query(ComplianceProgram)
            .filter(ComplianceProgram.tenant_id.in_(tenant_ids))
            .order_by(ComplianceProgram.id.desc())
            .limit(800)
            .all()
        )
        for item in compliance_programs:
            framework_name = item.framework.name if item.framework else ""
            add(
                _catalog_entry(
                    source_type="compliance_program",
                    source_id=str(item.id),
                    title=item.name or f"Compliance Program {item.id}",
                    description=_merge_text(framework_name, item.status, item.description)[:900],
                    updated_at=_to_iso(item.target_date or item.start_date),
                    tenant_id=item.tenant_id,
                )
            )

        compliance_assessments = (
            db.query(GRCComplianceAssessment, ComplianceProgram)
            .join(ComplianceProgram, GRCComplianceAssessment.program_id == ComplianceProgram.id)
            .filter(ComplianceProgram.tenant_id.in_(tenant_ids))
            .order_by(GRCComplianceAssessment.id.desc())
            .limit(1600)
            .all()
        )
        for assessment, program in compliance_assessments:
            control_name = assessment.normalized_control.name if assessment.normalized_control else ""
            maturity_label = (
                str(assessment.maturity_level)
                if assessment.maturity_level is not None
                else "n/a"
            )
            add(
                _catalog_entry(
                    source_type="grc_compliance_assessment",
                    source_id=str(assessment.id),
                    title=f"{program.name} - {control_name or f'Assessment {assessment.id}'}",
                    description=_merge_text(assessment.status, assessment.notes, f"Maturity: {maturity_label}")[:900],
                    updated_at=_to_iso(assessment.assessed_at),
                    tenant_id=program.tenant_id,
                )
            )

        try:
            framework_risk_assessments = (
                db.query(FrameworkRiskAssessment)
                .filter(FrameworkRiskAssessment.tenant_id.in_(tenant_ids))
                .order_by(FrameworkRiskAssessment.updated_at.desc())
                .limit(500)
                .all()
            )
        except Exception as exc:
            logger.warning("Skipping framework risk assessments for vector catalog due schema mismatch: %s", exc)
            framework_risk_assessments = []
        for item in framework_risk_assessments:
            framework_name = ""
            if item.uploaded_framework and item.uploaded_framework.name:
                framework_name = item.uploaded_framework.name
            elif item.framework and item.framework.name:
                framework_name = item.framework.name
            add(
                _catalog_entry(
                    source_type="framework_risk_assessment",
                    source_id=str(item.id),
                    title=item.name or f"Framework Risk Assessment {item.id}",
                    description=_merge_text(framework_name, item.status, item.description)[:800],
                    updated_at=_to_iso(item.updated_at or item.created_at),
                    tenant_id=item.tenant_id,
                )
            )

        try:
            compliance_docs = (
                db.query(ComplianceAssessmentDocument)
                .filter(ComplianceAssessmentDocument.tenant_id.in_(tenant_ids))
                .order_by(ComplianceAssessmentDocument.updated_at.desc())
                .limit(700)
                .all()
            )
        except Exception as exc:
            logger.warning("Skipping compliance assessment documents for vector catalog due schema mismatch: %s", exc)
            compliance_docs = []
        for item in compliance_docs:
            add(
                _catalog_entry(
                    source_type="compliance_assessment_document",
                    source_id=str(item.id),
                    title=item.name or f"Compliance Assessment {item.id}",
                    description=_merge_text(item.assessment_type, item.status, item.source, item.notes)[:800],
                    updated_at=_to_iso(item.updated_at or item.created_at),
                    tenant_id=item.tenant_id,
                )
            )

        evidence_rows = (
            db.query(Evidence)
            .filter(Evidence.tenant_id.in_(tenant_ids))
            .order_by(Evidence.uploaded_at.desc())
            .limit(1200)
            .all()
        )
        for item in evidence_rows:
            add(
                _catalog_entry(
                    source_type="evidence_file",
                    source_id=str(item.id),
                    title=item.name or item.file_name or f"Evidence {item.id}",
                    description=_merge_text(item.description, item.evidence_type, item.content_summary)[:800],
                    updated_at=_to_iso(item.uploaded_at),
                    tenant_id=item.tenant_id,
                )
            )

    uploaded_frameworks_query = (
        db.query(UploadedFramework)
        .filter(UploadedFramework.is_active == True)
        .order_by(UploadedFramework.updated_at.desc())
    )
    if tenant_ids:
        uploaded_frameworks_query = uploaded_frameworks_query.filter(
            or_(
                UploadedFramework.tenant_id.in_(tenant_ids),
                UploadedFramework.is_shared == True,
                UploadedFramework.tenant_id.is_(None),
            )
        )
    uploaded_frameworks = uploaded_frameworks_query.limit(500).all()
    for item in uploaded_frameworks:
        add(
            _catalog_entry(
                source_type="uploaded_framework",
                source_id=str(item.id),
                title=item.name or f"Uploaded Framework {item.id}",
                description=_merge_text(
                    item.framework_type,
                    item.classification,
                    item.version,
                    item.description,
                    item.framework_scope,
                )[:900],
                updated_at=_to_iso(item.updated_at or item.created_at),
                tenant_id=item.tenant_id,
            )
        )

    frameworks = (
        db.query(Framework)
        .filter(Framework.is_active == True)
        .order_by(Framework.name.asc())
        .limit(200)
        .all()
    )
    for item in frameworks:
        add(
            _catalog_entry(
                source_type="seeded_framework",
                source_id=str(item.id),
                title=item.name or f"Framework {item.id}",
                description=_merge_text(item.short_code, item.version, item.regulator, item.description)[:900],
                updated_at="",
                tenant_id=None,
            )
        )

    for item in uploaded_files or []:
        file_id = str(item.get("id") or uuid4().hex)
        title = str(item.get("filename") or "Uploaded File")
        add(
            _catalog_entry(
                source_type="chat_upload",
                source_id=file_id,
                title=title,
                description=str(item.get("excerpt") or "")[:700],
                updated_at=str(item.get("uploaded_at") or ""),
                tenant_id=None,
            )
        )

    return entries


def _title_match_score(question: str, title: str) -> float:
    if not title:
        return 0.0
    q_norm = normalize_question(question.lower())
    t_norm = normalize_question(title.lower())
    if not t_norm:
        return 0.0
    if t_norm in q_norm:
        return 1.0

    q_tokens = {token for token in re.findall(r"[a-z0-9]{3,}", q_norm)}
    t_tokens = {token for token in re.findall(r"[a-z0-9]{3,}", t_norm)}
    if not q_tokens or not t_tokens:
        return 0.0
    overlap = len(q_tokens.intersection(t_tokens))
    if overlap == 0:
        return 0.0
    return overlap / max(1, len(t_tokens))


def _select_vector_catalog_matches(
    question: str,
    catalog: List[Dict[str, Any]],
    max_results: int = VECTOR_MATCH_LIMIT,
) -> List[Dict[str, Any]]:
    if not catalog:
        return []

    normalized_question = normalize_question((question or "").lower())
    has_vector_intent = any(term in normalized_question for term in VECTOR_ROUTE_TERMS)

    scored: List[Tuple[float, Dict[str, Any]]] = []
    for entry in catalog:
        title_score = _title_match_score(question, entry.get("title", ""))
        desc_score = _title_match_score(question, entry.get("description", ""))
        score = max(title_score, desc_score * 0.65)
        if has_vector_intent and entry.get("source_type") in {
            "governance_document",
            "governance_document_version",
            "committee_charter",
            "uploaded_framework",
            "seeded_framework",
            "seeded_framework_control",
            "compliance_assessment_document",
            "compliance_program",
            "grc_compliance_assessment",
            "risk_assessment",
            "framework_risk_assessment",
            "chat_upload",
            "evidence_file",
            "policy_statement",
            "policy_statement_version",
        }:
            score = max(score, 0.12)
        if score > 0:
            scored.append((score, entry))

    if not scored:
        return []

    scored.sort(key=lambda item: item[0], reverse=True)
    return [entry for _, entry in scored[:max_results]]


def _has_vector_semantic_intent(question: str) -> bool:
    normalized = normalize_question((question or "").lower())
    return any(term in normalized for term in VECTOR_SEMANTIC_TERMS)


def _entry_tenants(entry: Dict[str, Any], tenant_ids: List[int]) -> List[int]:
    tenant_id = entry.get("tenant_id")
    if tenant_id:
        return [int(tenant_id)]
    return list(tenant_ids)


def _collect_documents_for_matches(
    db: Session,
    tenant_ids: List[int],
    matches: List[Dict[str, Any]],
    uploaded_files: Optional[List[Dict[str, Any]]] = None,
) -> List[Any]:
    if IndexedSourceDocument is None:
        return []

    documents: List[Any] = []
    upload_lookup = {str(item.get("id")): item for item in (uploaded_files or [])}

    for entry in matches:
        source_type = entry.get("source_type")
        source_id = str(entry.get("source_id"))
        tenants = _entry_tenants(entry, tenant_ids)
        if not tenants:
            continue

        if source_type == "governance_document":
            row = db.query(GovernanceDocument).filter(GovernanceDocument.id == int(source_id)).first()
            if not row:
                continue
            file_text = extract_text_from_path(row.file_path) if extract_text_from_path else ""
            base_text = _merge_text(
                row.title,
                row.description,
                row.content,
                file_text,
            )
            if not base_text.strip():
                continue
            for tenant_id in tenants:
                documents.append(
                    IndexedSourceDocument(
                        tenant_id=tenant_id,
                        source_type="governance_document",
                        source_id=source_id,
                        title=row.title or f"Governance Document {source_id}",
                        description=_merge_text(row.doc_type, row.status),
                        text=base_text,
                        updated_at=_to_iso(row.updated_at or row.created_at),
                        metadata={"doc_type": row.doc_type or "", "status": row.status or ""},
                    )
                )

        elif source_type == "governance_document_version":
            row = (
                db.query(GovernanceDocumentVersion, GovernanceDocument)
                .join(GovernanceDocument, GovernanceDocumentVersion.document_id == GovernanceDocument.id)
                .filter(GovernanceDocumentVersion.id == int(source_id))
                .first()
            )
            if not row:
                continue
            version, document = row
            file_text = extract_text_from_path(version.file_path) if extract_text_from_path else ""
            base_text = _merge_text(
                document.title,
                f"Version {version.version_number}",
                version.title,
                version.content,
                version.change_summary,
                version.change_reason,
                file_text,
            )
            if not base_text.strip():
                continue
            for tenant_id in tenants:
                documents.append(
                    IndexedSourceDocument(
                        tenant_id=tenant_id,
                        source_type="governance_document_version",
                        source_id=source_id,
                        title=version.title or document.title or f"Document Version {source_id}",
                        description=_merge_text(document.doc_type, version.status, version.version_number),
                        text=base_text,
                        updated_at=_to_iso(version.created_at),
                        metadata={"document_id": str(document.id), "version": version.version_number or ""},
                    )
                )

        elif source_type == "policy_statement":
            row = db.query(PolicyStatement).filter(PolicyStatement.id == int(source_id)).first()
            if not row:
                continue
            version_rows = (
                db.query(PolicyStatementVersion)
                .filter(PolicyStatementVersion.statement_id == row.id)
                .order_by(PolicyStatementVersion.version_number.desc())
                .limit(12)
                .all()
            )
            version_text = "\n".join(
                _merge_text(
                    f"v{version.version_number}",
                    version.statement_text,
                    version.statement_summary,
                    version.change_reason,
                )
                for version in version_rows
            )
            base_text = _merge_text(
                row.statement_code,
                row.statement_text,
                row.statement_summary,
                row.source_section,
                row.category,
                row.sub_category,
                _json_text(row.ai_extracted_keywords),
                _json_text(row.ai_suggested_controls),
                version_text,
            )
            if not base_text.strip():
                continue
            title = row.statement_code or f"Policy Statement {row.id}"
            for tenant_id in tenants:
                documents.append(
                    IndexedSourceDocument(
                        tenant_id=tenant_id,
                        source_type="policy_statement",
                        source_id=source_id,
                        title=title,
                        description=_merge_text(row.category, row.sub_category, row.status),
                        text=base_text,
                        updated_at=_to_iso(row.updated_at or row.created_at),
                        metadata={"document_id": str(row.document_id), "status": row.status or ""},
                    )
                )

        elif source_type == "policy_statement_version":
            row = db.query(PolicyStatementVersion).filter(PolicyStatementVersion.id == int(source_id)).first()
            if not row:
                continue
            base_text = _merge_text(
                f"Policy Statement {row.statement_id}",
                f"Version {row.version_number}",
                row.statement_text,
                row.statement_summary,
                row.category,
                row.sub_category,
                row.change_reason,
            )
            if not base_text.strip():
                continue
            for tenant_id in tenants:
                documents.append(
                    IndexedSourceDocument(
                        tenant_id=tenant_id,
                        source_type="policy_statement_version",
                        source_id=source_id,
                        title=f"Policy Statement {row.statement_id} v{row.version_number}",
                        description=_merge_text(row.category, row.sub_category, row.status),
                        text=base_text,
                        updated_at=_to_iso(row.changed_at),
                        metadata={"statement_id": str(row.statement_id), "status": row.status or ""},
                    )
                )

        elif source_type == "committee_charter":
            row = db.query(CommitteeCharter).filter(CommitteeCharter.id == int(source_id)).first()
            if not row:
                continue
            committee_name = row.committee.name if row.committee else ""
            file_text = extract_text_from_path(row.file_path) if extract_text_from_path else ""
            base_text = _merge_text(row.title, committee_name, row.content, file_text)
            if not base_text.strip():
                continue
            for tenant_id in tenants:
                documents.append(
                    IndexedSourceDocument(
                        tenant_id=tenant_id,
                        source_type="committee_charter",
                        source_id=source_id,
                        title=row.title or committee_name or f"Committee Charter {source_id}",
                        description=_merge_text(committee_name, row.version, row.status),
                        text=base_text,
                        updated_at=_to_iso(row.created_at),
                        metadata={"committee": committee_name, "status": row.status or ""},
                    )
                )

        elif source_type == "uploaded_framework":
            row = db.query(UploadedFramework).filter(UploadedFramework.id == int(source_id)).first()
            if not row:
                continue
            file_text = extract_text_from_path(row.file_path) if extract_text_from_path else ""
            base_text = _merge_text(
                row.name,
                row.description,
                row.framework_purpose,
                row.framework_scope,
                row.classification_reasoning,
                _json_text(row.framework_objectives),
                _json_text(row.adoption_approach),
                _json_text(row.hierarchy_structure),
                file_text,
            )
            if base_text.strip():
                for tenant_id in tenants:
                    documents.append(
                        IndexedSourceDocument(
                            tenant_id=tenant_id,
                            source_type="uploaded_framework",
                            source_id=source_id,
                            title=row.name or f"Uploaded Framework {source_id}",
                            description=_merge_text(row.framework_type, row.classification, row.version),
                            text=base_text,
                            updated_at=_to_iso(row.updated_at or row.created_at),
                            metadata={"framework_type": row.framework_type or "", "version": row.version or ""},
                        )
                    )

            controls = (
                db.query(ParsedFrameworkControl)
                .filter(ParsedFrameworkControl.uploaded_framework_id == row.id)
                .order_by(ParsedFrameworkControl.control_id.asc())
                .limit(max(400, VECTOR_PARSED_CONTROL_LIMIT))
                .all()
            )
            for control in controls:
                evidence_requirements = (
                    db.query(ControlEvidenceRequirement)
                    .filter(
                        ControlEvidenceRequirement.framework_id == row.id,
                        ControlEvidenceRequirement.parsed_control_id == control.id,
                    )
                    .order_by(ControlEvidenceRequirement.display_order.asc(), ControlEvidenceRequirement.id.asc())
                    .limit(30)
                    .all()
                )
                evidence_requirement_text = "\n".join(
                    _merge_text(
                        req.evidence_title,
                        req.evidence_description,
                        req.evidence_type,
                        req.evidence_format,
                        _json_text(req.exact_requirements),
                        _json_text(req.acceptance_criteria),
                        req.collection_guidance,
                        req.ai_reasoning,
                    )
                    for req in evidence_requirements
                )
                control_text = _merge_text(
                    control.control_id,
                    control.title,
                    control.description,
                    control.full_text,
                    control.domain,
                    control.category,
                    _json_text(control.evidence_requirements),
                    evidence_requirement_text,
                )
                if not control_text.strip():
                    continue
                for tenant_id in tenants:
                    documents.append(
                        IndexedSourceDocument(
                            tenant_id=tenant_id,
                            source_type="parsed_framework_control",
                            source_id=str(control.id),
                            title=f"{row.name} - {control.control_id} {control.title}".strip(),
                            description=_merge_text(control.domain, control.category),
                            text=control_text,
                            updated_at=_to_iso(control.updated_at or control.created_at),
                            metadata={
                                "framework_name": row.name or "",
                                "control_id": control.control_id or "",
                                "domain": control.domain or "",
                            },
                        )
                    )

        elif source_type == "seeded_framework":
            row = db.query(Framework).filter(Framework.id == int(source_id)).first()
            if not row:
                continue
            controls = (
                db.query(FrameworkControl, ControlObjective, FrameworkDomain)
                .join(ControlObjective, FrameworkControl.objective_id == ControlObjective.id)
                .join(FrameworkDomain, ControlObjective.domain_id == FrameworkDomain.id)
                .filter(FrameworkDomain.framework_id == row.id)
                .order_by(FrameworkControl.code.asc())
                .limit(max(600, VECTOR_SEEDED_CONTROL_LIMIT))
                .all()
            )
            framework_summary = _merge_text(
                row.name,
                row.description,
                row.regulator,
                row.jurisdiction,
                row.version,
            )
            if framework_summary.strip():
                for tenant_id in tenants:
                    documents.append(
                        IndexedSourceDocument(
                            tenant_id=tenant_id,
                            source_type="seeded_framework",
                            source_id=source_id,
                            title=row.name or f"Framework {source_id}",
                            description=_merge_text(row.short_code, row.version, row.regulator),
                            text=framework_summary,
                            updated_at="",
                            metadata={"short_code": row.short_code or "", "version": row.version or ""},
                        )
                    )

            control_id_list = [control.id for control, _, _ in controls]
            sub_controls_by_control: Dict[int, List[FrameworkSubControl]] = {}
            curated_by_control: Dict[int, List[CuratedEvidenceItem]] = {}
            if control_id_list:
                sub_controls = (
                    db.query(FrameworkSubControl)
                    .filter(FrameworkSubControl.control_id.in_(control_id_list))
                    .order_by(FrameworkSubControl.control_id.asc(), FrameworkSubControl.order.asc(), FrameworkSubControl.id.asc())
                    .all()
                )
                for sub_control in sub_controls:
                    sub_controls_by_control.setdefault(int(sub_control.control_id), []).append(sub_control)

                curated_rows = (
                    db.query(CuratedEvidenceItem)
                    .filter(
                        or_(
                            CuratedEvidenceItem.framework_control_id.in_(control_id_list),
                            CuratedEvidenceItem.sub_control_id.in_([sub.id for sub in sub_controls]) if sub_controls else False,
                        )
                    )
                    .order_by(CuratedEvidenceItem.id.asc())
                    .all()
                )
                for item in curated_rows:
                    if item.framework_control_id:
                        curated_by_control.setdefault(int(item.framework_control_id), []).append(item)
                        continue
                    if item.sub_control_id:
                        parent = next((sub.control_id for sub in sub_controls if sub.id == item.sub_control_id), None)
                        if parent:
                            curated_by_control.setdefault(int(parent), []).append(item)

            for control, objective, domain in controls:
                related_sub_controls = sub_controls_by_control.get(int(control.id), [])
                related_curated = curated_by_control.get(int(control.id), [])
                sub_control_text = "\n".join(
                    _merge_text(
                        f"{sub.code}: {sub.name}",
                        sub.statement,
                        sub.description,
                        _json_text(sub.evidence_recommendations),
                    )
                    for sub in related_sub_controls
                )
                curated_text = "\n".join(
                    _merge_text(
                        item.title,
                        item.description,
                        item.artifact_type,
                        item.format_guidance,
                        item.frequency,
                    )
                    for item in related_curated
                )
                control_text = _merge_text(
                    f"{control.code}: {control.name}",
                    control.statement,
                    control.control_objective,
                    f"Domain: {domain.code} {domain.name}",
                    f"Objective: {objective.code} {objective.name}",
                    control.implementation_guidance,
                    control.testing_guidance,
                    f"Risk Category: {control.risk_category}" if control.risk_category else "",
                    f"Evidence Type: {control.evidence_type}" if control.evidence_type else "",
                    sub_control_text,
                    curated_text,
                )
                if not control_text.strip():
                    continue
                for tenant_id in tenants:
                    documents.append(
                        IndexedSourceDocument(
                            tenant_id=tenant_id,
                            source_type="seeded_framework_control",
                            source_id=str(control.id),
                            title=f"{row.name} - {control.code} {control.name}".strip(),
                            description=_merge_text(domain.name, objective.name, row.short_code),
                            text=control_text,
                            updated_at="",
                            metadata={
                                "framework_name": row.name or "",
                                "framework_code": row.short_code or "",
                                "control_code": control.code or "",
                                "domain": domain.name or "",
                                "objective": objective.name or "",
                            },
                        )
                    )

        elif source_type == "risk_assessment":
            row = db.query(RiskAssessment).filter(RiskAssessment.id == int(source_id)).first()
            if not row:
                continue
            risk_rows = (
                db.query(RiskAssessmentRisk, Risk)
                .join(Risk, RiskAssessmentRisk.risk_id == Risk.id)
                .filter(RiskAssessmentRisk.assessment_id == row.id)
                .limit(160)
                .all()
            )
            risk_summary = "\n".join(
                _merge_text(
                    f"Risk: {getattr(risk_record, 'title', None) or getattr(risk_record, 'name', None) or f'Risk {risk_record.id}'}",
                    f"Treatment: {assessment_risk.treatment_decision or 'n/a'}",
                    f"Rating: {assessment_risk.risk_rating or 'n/a'}",
                    assessment_risk.rationale,
                    assessment_risk.notes,
                )
                for assessment_risk, risk_record in risk_rows
            )
            text_blob = _merge_text(
                row.name,
                row.description,
                row.scope,
                row.methodology,
                row.notes,
                risk_summary,
            )
            if not text_blob.strip():
                continue
            for tenant_id in tenants:
                documents.append(
                    IndexedSourceDocument(
                        tenant_id=tenant_id,
                        source_type="risk_assessment",
                        source_id=source_id,
                        title=row.name or f"Risk Assessment {source_id}",
                        description=_merge_text(row.assessment_type, row.status),
                        text=text_blob,
                        updated_at=_to_iso(row.updated_at or row.created_at),
                        metadata={"assessment_type": row.assessment_type or "", "status": row.status or ""},
                    )
                )

        elif source_type == "compliance_program":
            row = db.query(ComplianceProgram).filter(ComplianceProgram.id == int(source_id)).first()
            if not row:
                continue
            assessment_rows = (
                db.query(GRCComplianceAssessment)
                .filter(GRCComplianceAssessment.program_id == row.id)
                .order_by(GRCComplianceAssessment.id.asc())
                .limit(700)
                .all()
            )
            assessment_text = "\n".join(
                _merge_text(
                    f"Assessment #{assessment.id}",
                    f"Status: {assessment.status}",
                    f"Maturity Level: {assessment.maturity_level if assessment.maturity_level is not None else 'n/a'}",
                    assessment.notes,
                    assessment.normalized_control.name if assessment.normalized_control else "",
                    assessment.normalized_control.statement if assessment.normalized_control else "",
                )
                for assessment in assessment_rows
            )
            framework_name = row.framework.name if row.framework else ""
            text_blob = _merge_text(
                row.name,
                framework_name,
                row.description,
                row.status,
                assessment_text,
            )
            if not text_blob.strip():
                continue
            for tenant_id in tenants:
                documents.append(
                    IndexedSourceDocument(
                        tenant_id=tenant_id,
                        source_type="compliance_program",
                        source_id=source_id,
                        title=row.name or f"Compliance Program {source_id}",
                        description=_merge_text(framework_name, row.status),
                        text=text_blob,
                        updated_at=_to_iso(row.target_date or row.start_date),
                        metadata={"framework_name": framework_name, "status": row.status or ""},
                    )
                )

        elif source_type == "grc_compliance_assessment":
            row = db.query(GRCComplianceAssessment).filter(GRCComplianceAssessment.id == int(source_id)).first()
            if not row:
                continue
            control_name = row.normalized_control.name if row.normalized_control else ""
            control_statement = row.normalized_control.statement if row.normalized_control else ""
            program_name = row.program.name if row.program else ""
            text_blob = _merge_text(
                program_name,
                control_name,
                control_statement,
                row.notes,
                f"Status: {row.status}",
                f"Maturity Level: {row.maturity_level if row.maturity_level is not None else 'n/a'}",
            )
            if not text_blob.strip():
                continue
            for tenant_id in tenants:
                documents.append(
                    IndexedSourceDocument(
                        tenant_id=tenant_id,
                        source_type="grc_compliance_assessment",
                        source_id=source_id,
                        title=f"{program_name} - {control_name or f'Assessment {source_id}'}",
                        description=_merge_text(row.status, f"Maturity: {row.maturity_level if row.maturity_level is not None else 'n/a'}"),
                        text=text_blob,
                        updated_at=_to_iso(row.assessed_at),
                        metadata={"program_id": str(row.program_id), "status": row.status or ""},
                    )
                )

        elif source_type == "framework_risk_assessment":
            row = db.query(FrameworkRiskAssessment).filter(FrameworkRiskAssessment.id == int(source_id)).first()
            if not row:
                continue
            question_rows = (
                db.query(FrameworkRiskQuestion)
                .filter(FrameworkRiskQuestion.assessment_id == row.id)
                .order_by(FrameworkRiskQuestion.order_index.asc())
                .limit(260)
                .all()
            )
            questions_text = "\n".join(
                _merge_text(question.question_text, f"Status: {question.status}") for question in question_rows
            )
            framework_name = ""
            if row.uploaded_framework and row.uploaded_framework.name:
                framework_name = row.uploaded_framework.name
            elif row.framework and row.framework.name:
                framework_name = row.framework.name
            text_blob = _merge_text(
                row.name,
                row.description,
                framework_name,
                questions_text,
            )
            if not text_blob.strip():
                continue
            for tenant_id in tenants:
                documents.append(
                    IndexedSourceDocument(
                        tenant_id=tenant_id,
                        source_type="framework_risk_assessment",
                        source_id=source_id,
                        title=row.name or f"Framework Risk Assessment {source_id}",
                        description=_merge_text(framework_name, row.status),
                        text=text_blob,
                        updated_at=_to_iso(row.updated_at or row.created_at),
                        metadata={"framework_name": framework_name, "status": row.status or ""},
                    )
                )

        elif source_type == "compliance_assessment_document":
            row = db.query(ComplianceAssessmentDocument).filter(ComplianceAssessmentDocument.id == int(source_id)).first()
            if not row:
                continue
            item_rows = (
                db.query(ComplianceAssessmentDocumentItem)
                .filter(ComplianceAssessmentDocumentItem.assessment_id == row.id)
                .order_by(ComplianceAssessmentDocumentItem.id.asc())
                .limit(300)
                .all()
            )
            item_text = "\n".join(
                _merge_text(
                    item.item_number,
                    item.area_domain,
                    item.control_description,
                    f"Status: {item.compliance_status}",
                    item.gaps_identified,
                    item.proposed_solution,
                    item.remarks,
                )
                for item in item_rows
            )
            file_text = extract_text_from_path(row.file_path) if extract_text_from_path else ""
            text_blob = _merge_text(
                row.name,
                row.notes,
                row.assessment_type,
                row.source,
                _json_text(row.xlsx_data),
                item_text,
                file_text,
            )
            if not text_blob.strip():
                continue
            for tenant_id in tenants:
                documents.append(
                    IndexedSourceDocument(
                        tenant_id=tenant_id,
                        source_type="compliance_assessment_document",
                        source_id=source_id,
                        title=row.name or f"Compliance Assessment Document {source_id}",
                        description=_merge_text(row.assessment_type, row.status, row.source),
                        text=text_blob,
                        updated_at=_to_iso(row.updated_at or row.created_at),
                        metadata={"assessment_type": row.assessment_type or "", "status": row.status or ""},
                    )
                )

        elif source_type == "evidence_file":
            row = db.query(Evidence).filter(Evidence.id == int(source_id)).first()
            if not row:
                continue
            file_text = extract_text_from_path(row.file_path) if extract_text_from_path else ""
            text_blob = _merge_text(
                row.name,
                row.description,
                row.content_summary,
                row.ocr_content,
                file_text,
            )
            if not text_blob.strip():
                continue
            for tenant_id in tenants:
                documents.append(
                    IndexedSourceDocument(
                        tenant_id=tenant_id,
                        source_type="evidence_file",
                        source_id=source_id,
                        title=row.name or row.file_name or f"Evidence {source_id}",
                        description=_merge_text(row.evidence_type, row.status),
                        text=text_blob,
                        updated_at=_to_iso(row.uploaded_at),
                        metadata={"evidence_type": row.evidence_type or "", "status": row.status or ""},
                    )
                )

        elif source_type == "chat_upload":
            uploaded = upload_lookup.get(source_id)
            if not uploaded:
                continue
            upload_text = _merge_text(
                uploaded.get("excerpt"),
                extract_text_from_path(uploaded.get("path")) if extract_text_from_path else "",
            )
            if not upload_text.strip():
                continue
            for tenant_id in tenants:
                documents.append(
                    IndexedSourceDocument(
                        tenant_id=tenant_id,
                        source_type="chat_upload",
                        source_id=source_id,
                        title=str(uploaded.get("filename") or f"Chat Upload {source_id}"),
                        description="Session uploaded file",
                        text=upload_text,
                        updated_at=str(uploaded.get("uploaded_at") or ""),
                        metadata={"filename": str(uploaded.get("filename") or "")},
                    )
                )

    return documents


def _vector_doc_signature(doc: Any) -> str:
    raw = f"{doc.source_type}|{doc.source_id}|{doc.updated_at}|{len(doc.text)}|{doc.title}"
    return hashlib.sha1(raw.encode("utf-8")).hexdigest()


def _sync_docs_to_qdrant(service: Any, docs: List[Any], *, force: bool = False) -> int:
    if not docs:
        return 0
    to_upsert: List[Any] = []
    updated_cache: List[Tuple[int, str, str]] = []

    for doc in docs:
        cache_key = f"{doc.source_type}:{doc.source_id}"
        signature = _vector_doc_signature(doc)
        tenant_cache = VECTOR_INDEX_STATE.setdefault(doc.tenant_id, {})
        if not force and tenant_cache.get(cache_key) == signature:
            continue
        try:
            service.delete_source_points(doc.tenant_id, doc.source_type, str(doc.source_id))
        except Exception:
            # Best effort: continue with upsert.
            pass
        to_upsert.append(doc)
        updated_cache.append((doc.tenant_id, cache_key, signature))

    if not to_upsert:
        return 0

    indexed_points = service.upsert_documents(to_upsert)
    for tenant_id, cache_key, signature in updated_cache:
        VECTOR_INDEX_STATE.setdefault(tenant_id, {})[cache_key] = signature
    return indexed_points


def _index_chat_uploads_to_qdrant(
    *,
    service: Any,
    tenant_ids: List[int],
    uploaded_items: List[Dict[str, Any]],
) -> int:
    if not uploaded_items or not tenant_ids or IndexedSourceDocument is None:
        return 0

    docs: List[Any] = []
    for item in uploaded_items:
        source_id = str(item.get("id") or "")
        if not source_id:
            continue
        upload_text = _merge_text(
            item.get("excerpt"),
            extract_text_from_path(item.get("path")) if extract_text_from_path else "",
        )
        if not upload_text.strip():
            continue
        title = str(item.get("filename") or f"Chat Upload {source_id}")
        updated_at = str(item.get("uploaded_at") or "")
        for tenant_id in tenant_ids:
            docs.append(
                IndexedSourceDocument(
                    tenant_id=int(tenant_id),
                    source_type="chat_upload",
                    source_id=source_id,
                    title=title,
                    description="Session uploaded file",
                    text=upload_text,
                    updated_at=updated_at,
                    metadata={"filename": title},
                )
            )

    if not docs:
        return 0
    return _sync_docs_to_qdrant(service, docs, force=False)


def _query_qdrant_hits(
    service: Any,
    tenant_ids: List[int],
    question: str,
    *,
    limit: int = 8,
    source_types: Optional[Sequence[str]] = None,
) -> List[Dict[str, Any]]:
    all_hits: List[Dict[str, Any]] = []
    for tenant_id in tenant_ids:
        try:
            tenant_hits = service.search(
                tenant_id=tenant_id,
                query=question,
                limit=limit,
                source_types=source_types,
            )
            all_hits.extend(tenant_hits)
        except Exception as exc:
            logger.warning("Qdrant search failed for tenant %s: %s", tenant_id, exc)

    deduped: Dict[str, Dict[str, Any]] = {}
    for hit in all_hits:
        point_id = str(hit.get("id"))
        if point_id not in deduped or float(hit.get("score") or 0.0) > float(deduped[point_id].get("score") or 0.0):
            deduped[point_id] = hit

    sorted_hits = sorted(deduped.values(), key=lambda item: float(item.get("score") or 0.0), reverse=True)
    return [item for item in sorted_hits if float(item.get("score") or 0.0) >= VECTOR_SCORE_THRESHOLD][:limit]


def _vector_hits_to_sources(hits: List[Dict[str, Any]]) -> List[ChatSource]:
    sources: List[ChatSource] = []
    for index, hit in enumerate(hits, start=1):
        payload = hit.get("payload") or {}
        source_type = str(payload.get("source_type") or "vector_context")
        source_id = str(payload.get("source_id") or "")
        title = str(payload.get("title") or "")
        snippet = str(payload.get("snippet") or payload.get("text") or "")[:280]
        score_value = float(hit.get("score") or 0.0)
        sources.append(
            ChatSource(
                rank=index,
                entity_type=source_type,
                entity_id=source_id,
                framework_code=source_type.replace("_", " ").title(),
                control_code=source_id or None,
                control_name=title or None,
                relevance_score=score_value,
                snippet=snippet,
            )
        )
    return sources


def _answer_with_vector_context(
    question: str,
    hits: List[Dict[str, Any]],
    *,
    context_summary: str = "",
    uploaded_context: str = "",
) -> str:
    if not hits:
        return ""

    context_blocks: List[str] = []
    for index, hit in enumerate(hits[:10], start=1):
        payload = hit.get("payload") or {}
        title = str(payload.get("title") or "Untitled Source")
        source_type = str(payload.get("source_type") or "source")
        snippet = str(payload.get("text") or payload.get("snippet") or "")
        context_blocks.append(
            f"[Source {index}] {title} ({source_type})\n{snippet[:1800]}"
        )

    db_context = "Retrieved platform document context:\n\n" + "\n\n".join(context_blocks)
    return answer_grc_knowledge_question(
        question=question,
        context_summary=context_summary,
        uploaded_context=uploaded_context,
        db_context=db_context,
    )


def try_qdrant_rag_answer(
    *,
    question: str,
    request_mode: str,
    db: Session,
    tenant_ids: List[int],
    context_summary: str,
    uploaded_context: str,
    uploaded_files: Optional[List[Dict[str, Any]]] = None,
) -> Optional[Tuple[str, List[ChatSource]]]:
    if request_mode in {"framework_progress", "evidence_gaps"}:
        return None
    if not tenant_ids:
        return None

    service = get_qdrant_service()
    if not service or not getattr(service, "is_available", False):
        return None

    normalized = normalize_question((question or "").lower())
    has_vector_intent = any(term in normalized for term in VECTOR_ROUTE_TERMS)
    semantic_intent = _has_vector_semantic_intent(question)
    if request_mode == "database" and not semantic_intent:
        return None

    catalog = _build_vector_catalog(db, tenant_ids, uploaded_files)
    matches = _select_vector_catalog_matches(question, catalog)
    if not matches and not has_vector_intent and not semantic_intent:
        return None

    selected_entries = list(matches)
    if semantic_intent:
        selected_entries = list(catalog)
    elif not selected_entries and has_vector_intent:
        selected_entries = sorted(
            catalog,
            key=lambda entry: str(entry.get("updated_at") or ""),
            reverse=True,
        )[: min(VECTOR_CATALOG_FALLBACK_LIMIT, len(catalog))]

    if not selected_entries:
        return None

    if has_vector_intent and not semantic_intent:
        try:
            total_points = sum(service.count_points(tenant_id=tenant_id) for tenant_id in tenant_ids)
            if total_points == 0 and len(catalog) > len(selected_entries):
                selected_entries = list(catalog)
        except Exception:
            # Best effort only; continue with selected entries.
            pass

    docs = _collect_documents_for_matches(db, tenant_ids, selected_entries, uploaded_files)
    if not docs:
        return None
    if len(docs) > VECTOR_DOC_LIMIT_ON_ASK:
        docs = docs[:VECTOR_DOC_LIMIT_ON_ASK]

    try:
        indexed_count = _sync_docs_to_qdrant(service, docs, force=False)
        if indexed_count:
            logger.info("[QDRANT] Indexed %s point(s) for %s source document(s).", indexed_count, len(docs))
    except Exception as exc:
        logger.warning("Qdrant indexing failed: %s", exc)
        return None

    source_types = None if semantic_intent else sorted({str(doc.source_type) for doc in docs if getattr(doc, "source_type", "")})
    hits = _query_qdrant_hits(
        service=service,
        tenant_ids=tenant_ids,
        question=question,
        limit=8,
        source_types=source_types,
    )
    if not hits and source_types:
        hits = _query_qdrant_hits(
            service=service,
            tenant_ids=tenant_ids,
            question=question,
            limit=8,
            source_types=None,
        )
    if not hits:
        return None

    answer = _answer_with_vector_context(
        question=question,
        hits=hits,
        context_summary=context_summary,
        uploaded_context=uploaded_context,
    )
    if not answer.strip():
        return None

    sources = _vector_hits_to_sources(hits)
    return answer, sources


# ============================================================================
# ENDPOINTS
# ============================================================================

@router.post("/ask", response_model=ChatResponse)
async def ask_compliance_question(
    request: ChatRequest,
    response: Response,
    db: Session = Depends(get_db),
    current_user: GRCUser = Depends(require_auth)
):
    """
    Ask ANY question and get AI-powered answers from hybrid retrieval.
    
    **HYBRID INTELLIGENCE**
    - Live platform data questions route to SQL Agent (real-time database answers)
    - Document/policy/framework content questions route to Qdrant vector retrieval when relevant
    - Title-based routing decides whether vector search is needed
    - Conversation context is retained for follow-up questions
    
    **Example Questions:**
    - "List all critical vulnerabilities" [>] Direct SQL query
    - "Show me more details about the first one" [>] Uses context from previous query
    - "What does PCI DSS require for encryption?" [>] Query grc_framework_controls table
    - "How about ISO 27001?" [>] Understands reference to previous topic
    """
    # Add no-cache headers to prevent any caching
    response.headers["Cache-Control"] = "no-cache, no-store, must-revalidate, max-age=0"
    response.headers["Pragma"] = "no-cache"
    response.headers["Expires"] = "0"
    
    if not SQL_AGENT_ENABLED or not callable(generate_sql_query) or not callable(validate_sql):
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail="SQL Agent service is not available."
        )

    # === Deterministic handling for framework progress / active journey overview ===
    normalized_question = request.message.lower()
    if (
        "framework" in normalized_question
        and ("active" in normalized_question or "progress" in normalized_question)
        and ("journey" in normalized_question or "certification" in normalized_question or "progress" in normalized_question)
    ):
        tenant_ids = get_user_tenants(current_user, db)
        answer = summarize_framework_progress(db, tenant_ids)

        return ChatResponse(
            answer=answer,
            sources=[],
            framework_filter="FRAMEWORK_PROGRESS_OVERVIEW",
            timestamp=datetime.utcnow().isoformat(),
            has_more=False,
            total_count=0,
            current_offset=0
        )
    
    # Resolve session-specific context and hydrate LangChain history
    session_id = resolve_chat_session_id(request.session_id, current_user.id)
    recent_history = request.history or get_recent_session_log(current_user.id, session_id)
    context_summary = build_context_summary(current_user.id, session_id, recent_history)
    uploaded_files = load_session_uploaded_files(current_user.id, session_id)
    uploaded_context = build_uploaded_context(current_user.id, session_id)

    logger.info(f"\n{'='*60}\n🤖 NEW QUESTION: {request.message}\nSession: {session_id}\nHistory: {len(recent_history)} messages\nFramework filter: {request.framework}\n{'='*60}")

    try:
        request_mode = classify_request_mode(request.message, has_uploaded_files=bool(uploaded_files))

        if request_mode == "blocked_injection":
            answer = (
                "This request contains patterns that ComplyChat cannot process. "
                "Please ask a genuine GRC question about compliance frameworks, risks, controls, policies, or governance."
            )
            store_chat_exchange(current_user.id, session_id, request.message, answer, offset=request.offset)
            return ChatResponse(
                answer=answer,
                sources=[],
                framework_filter=request.framework,
                timestamp=datetime.utcnow().isoformat(),
                has_more=False,
                total_count=0,
                current_offset=0
            )

        if request_mode == "blocked_harmful":
            answer = (
                "ComplyChat cannot assist with requests related to bypassing, defeating, or circumventing "
                "security controls, audit trails, or compliance mechanisms. "
                "If you have a legitimate security testing requirement, please work through your authorized "
                "security team and follow your organization's change management and approval process."
            )
            store_chat_exchange(current_user.id, session_id, request.message, answer, offset=request.offset)
            return ChatResponse(
                answer=answer,
                sources=[],
                framework_filter=request.framework,
                timestamp=datetime.utcnow().isoformat(),
                has_more=False,
                total_count=0,
                current_offset=0
            )

        if request_mode == "off_topic":
            answer = (
                "ComplyChat is limited to GRC topics only. "
                "Please ask about compliance frameworks, ISO, PCI DSS, SBP, governance, evidence, risks, vulnerabilities, assets, or related organizational controls."
            )
            store_chat_exchange(current_user.id, session_id, request.message, answer, offset=request.offset)
            return ChatResponse(
                answer=answer,
                sources=[],
                framework_filter=request.framework,
                timestamp=datetime.utcnow().isoformat(),
                has_more=False,
                total_count=0,
                current_offset=0
            )

        if request_mode == "deprecated_audit":
            answer = (
                "Audit Management is no longer available in this system. "
                "I can still help with governance, policies, compliance, risks, evidence, certifications, and vulnerabilities."
            )
            store_chat_exchange(current_user.id, session_id, request.message, answer, offset=request.offset, is_error=False)
            return ChatResponse(
                answer=answer,
                sources=[],
                framework_filter=request.framework,
                timestamp=datetime.utcnow().isoformat(),
                has_more=False,
                total_count=0,
                current_offset=0
            )

        tenant_ids = get_user_tenants(current_user, db)

        if _is_policy_count_question(request.message):
            answer = _policy_count_answer(db, tenant_ids)
            store_chat_exchange(current_user.id, session_id, request.message, answer, offset=request.offset)
            return ChatResponse(
                answer=answer,
                sources=[],
                framework_filter="DB_POLICY_COUNT",
                timestamp=datetime.utcnow().isoformat(),
                has_more=False,
                total_count=0,
                current_offset=0
            )

        if request_mode == "framework_progress":
            answer = summarize_framework_progress(db, tenant_ids)
            store_chat_exchange(current_user.id, session_id, request.message, answer, offset=request.offset)
            return ChatResponse(
                answer=answer,
                sources=[],
                framework_filter=request.framework,
                timestamp=datetime.utcnow().isoformat(),
                has_more=False,
                total_count=0,
                current_offset=0
            )

        if request_mode == "evidence_gaps":
            answer = summarize_evidence_gaps(db, tenant_ids)
            store_chat_exchange(current_user.id, session_id, request.message, answer, offset=request.offset)
            return ChatResponse(
                answer=answer,
                sources=[],
                framework_filter=request.framework,
                timestamp=datetime.utcnow().isoformat(),
                has_more=False,
                total_count=0,
                current_offset=0
            )

        with usage_scope(
            module_key="complychat",
            feature_key="vector_rag",
            actor_user_id=current_user.id,
            actor_username=current_user.username,
        ):
            vector_rag_result = try_qdrant_rag_answer(
                question=request.message,
                request_mode=request_mode,
                db=db,
                tenant_ids=tenant_ids,
                context_summary=context_summary,
                uploaded_context=uploaded_context,
                uploaded_files=uploaded_files,
            )
        if vector_rag_result is not None:
            answer, vector_sources = vector_rag_result
            store_chat_exchange(
                current_user.id,
                session_id,
                request.message,
                answer,
                offset=request.offset,
                sources=[source.model_dump() for source in vector_sources],
            )
            return ChatResponse(
                answer=answer,
                sources=vector_sources if request.include_sources else [],
                framework_filter="QDRANT_VECTOR_RAG",
                timestamp=datetime.utcnow().isoformat(),
                has_more=False,
                total_count=len(vector_sources),
                current_offset=0
            )

        if request_mode == "file_analysis":
            logger.info(f"[FILE-ANALYSIS] Sending {len(uploaded_files)} file(s) + question to LLM: {request.message}")
            with usage_scope(
                module_key="complychat",
                feature_key="file_analysis",
                actor_user_id=current_user.id,
                actor_username=current_user.username,
            ):
                answer = analyze_uploaded_files_with_llm(request.message, uploaded_files, context_summary)
            store_chat_exchange(current_user.id, session_id, request.message, answer, offset=request.offset)
            return ChatResponse(
                answer=answer,
                sources=[],
                framework_filter=request.framework,
                timestamp=datetime.utcnow().isoformat(),
                has_more=False,
                total_count=0,
                current_offset=0
            )

        if request_mode == "grc_guidance":
            # LLM-driven answer — use DB-augmented context for framework/knowledge questions
            logger.info(f"[LLM-GUIDANCE] Routing to DB-augmented LLM guidance for: {request.message}")
            tenant_ids = get_user_tenants(current_user, db)
            with usage_scope(
                module_key="complychat",
                feature_key="knowledge_guidance",
                actor_user_id=current_user.id,
                actor_username=current_user.username,
            ):
                answer = answer_with_db_context(request.message, db, tenant_ids, context_summary, uploaded_context)
            store_chat_exchange(current_user.id, session_id, request.message, answer, offset=request.offset)
            return ChatResponse(
                answer=answer,
                sources=[],
                framework_filter=request.framework,
                timestamp=datetime.utcnow().isoformat(),
                has_more=False,
                total_count=0,
                current_offset=0
            )

        # Enhance question with context if user references previous conversation or uploaded files
        enhanced_question = request.message
        reference_words = ['that', 'it', 'them', 'those', 'this', 'first', 'second', 'last', 'previous', 'above']
        if any(word in request.message.lower() for word in reference_words) and context_summary:
            enhanced_question = f"{request.message}\n\n{context_summary}"
            logger.info(f"💡 Enhanced question with context: {enhanced_question}")
        if uploaded_context:
            enhanced_question = f"{enhanced_question}\n\n{uploaded_context}"
        
        tenant_ids = get_user_tenants(current_user, db)

        # 🤖 STEP 1: Generate SQL query from natural language
        logger.info("[STATS] Generating SQL query from question...")
        with usage_scope(
            module_key="complychat",
            feature_key="sql_generation",
            actor_user_id=current_user.id,
            actor_username=current_user.username,
        ):
            sql_result = generate_sql_query(enhanced_question, language="en", limit=request.limit, offset=request.offset)
        
        if not sql_result.get('sql') or not validate_sql(sql_result['sql']):
            # No valid SQL — return a deterministic, non-fabricated response.
            logger.info("[SAFE-FAIL] No valid SQL generated; returning deterministic response")
            answer = _sql_failure_answer(request.message)
            store_chat_exchange(current_user.id, session_id, request.message, answer, offset=request.offset)
            return ChatResponse(
                answer=answer,
                sources=[],
                framework_filter=request.framework,
                timestamp=datetime.utcnow().isoformat(),
                has_more=False,
                total_count=0,
                current_offset=0
            )
        
        # [SEARCH] STEP 2: Validate columns BEFORE execution
        sql_query = sql_result['sql']
        
        # Validate that all columns exist in database
        validation = validate_columns_in_sql(sql_query)
        if not validation['valid']:
            logger.warning(f"[WARN]️ Column validation failed: {validation['errors']}")
            
            # [REFRESH] SMART RETRY: Auto-fetch real schema and regenerate query
            logger.info("[REFRESH] ATTEMPTING SMART RETRY: Fetching actual table schemas...")
            
            # Extract table names from failed query
            import re
            table_pattern = r'FROM\s+(\w+)|JOIN\s+(\w+)'
            tables = set()
            for match in re.finditer(table_pattern, sql_query, re.IGNORECASE):
                table_name = match.group(1) or match.group(2)
                if table_name:
                    tables.add(table_name)
            
            logger.info(f"[STATS] Tables in failed query: {tables}")
            
            # Fetch real schemas for these tables
            updated_schema_info = []
            for table in tables:
                try:
                    schema_text = fetch_table_schema_from_db(table)
                    if schema_text:
                        updated_schema_info.append(schema_text)
                        logger.info(f"  [YES] Fetched schema for {table}")
                except Exception as e:
                    logger.warning(f"  [WARN]️ Could not fetch schema for {table}: {e}")
            
            if updated_schema_info:
                # Retry query generation with real schema
                logger.info("[REFRESH] Retrying query generation with actual schema...")
                retry_prompt = f"""
PREVIOUS QUERY FAILED - Column names were incorrect.

ACTUAL TABLE SCHEMAS (from database):
{chr(10).join(updated_schema_info)}

ORIGINAL QUESTION: {request.message}

Generate new SQL using ONLY the column names listed above.
"""
                
                with usage_scope(
                    module_key="complychat",
                    feature_key="sql_generation_retry",
                    actor_user_id=current_user.id,
                    actor_username=current_user.username,
                ):
                    retry_result = generate_sql_query(
                        retry_prompt,
                        language="en",
                        offset=request.offset,
                        limit=request.limit
                    )
                
                if retry_result and retry_result.get('sql'):
                    logger.info(f"[YES] RETRY SUCCESSFUL: New SQL generated")
                    sql_query = retry_result['sql']
                    sql_result = retry_result
                    # Re-validate new query
                    validation = validate_columns_in_sql(sql_query)
                    if not validation['valid']:
                        logger.error("[FAIL] RETRY FAILED: New query still has column errors")
                        answer = _sql_failure_answer(request.message)
                        store_chat_exchange(current_user.id, session_id, request.message, answer, offset=request.offset)
                        return ChatResponse(
                            answer=answer,
                            sources=[],
                            framework_filter=request.framework,
                            timestamp=datetime.utcnow().isoformat(),
                            has_more=False,
                            total_count=0,
                            current_offset=0
                        )
                else:
                    logger.error("[FAIL] RETRY FAILED: Could not regenerate query")
            
            # If retry failed or no schema found, return deterministic failure.
            if not validation['valid']:
                logger.error("[FAIL] Final validation failed")
                answer = _sql_failure_answer(request.message)
                store_chat_exchange(current_user.id, session_id, request.message, answer, offset=request.offset)
                return ChatResponse(
                    answer=answer,
                    sources=[],
                    framework_filter=request.framework,
                    timestamp=datetime.utcnow().isoformat(),
                    has_more=False,
                    total_count=0,
                    current_offset=0
                )
        
        # [SEARCH] STEP 3: Execute SQL query
        logger.info(f"[SEARCH] Executing SQL: {sql_query[:200]}...")
        
        try:
            # Get total count first (remove LIMIT/OFFSET for count query)
            count_query = sql_query.upper()
            if 'LIMIT' in count_query:
                count_query = sql_query[:sql_query.upper().rfind('LIMIT')].strip()
            count_sql = f"SELECT COUNT(*) as total FROM ({count_query}) as count_subquery"
            
            try:
                count_result = db.execute(text(count_sql))
                total_count = count_result.scalar() or 0
            except Exception as count_err:
                logger.warning(f"[WARN]️ Could not get total count: {count_err}")
                # Rollback to clear failed transaction
                try:
                    db.rollback()
                except:
                    pass
                total_count = 0
            
            # Execute paginated query
            result = db.execute(text(sql_query))
            rows = result.fetchall()
            
            # Convert rows to dicts
            data_list = [dict(row._mapping) for row in rows]
            logger.info(f"[YES] SQL returned {len(data_list)} rows (Total: {total_count}, Offset: {request.offset})")
            
            # Check if more results available
            has_more = (request.offset + len(data_list)) < total_count
            
            # [STYLE] STEP 3: Format results using AI
            # Add pagination context to formatting
            pagination_note = ""
            if total_count > request.limit:
                pagination_note = f"\n\n*Showing {request.offset + 1}-{request.offset + len(data_list)} of {total_count} total results*"

            if not data_list:
                logger.info("[SAFE-EMPTY] SQL returned 0 rows")
                answer = _no_platform_data_answer(request.message)
            else:
                with usage_scope(
                    module_key="complychat",
                    feature_key="sql_result_formatting",
                    actor_user_id=current_user.id,
                    actor_username=current_user.username,
                ):
                    answer = format_query_results(data_list, request.message, sql_query, language="en") + pagination_note
            
            # 💾 STEP 4: Save to conversation history
            store_chat_exchange(
                current_user.id,
                session_id,
                request.message,
                answer[:500],
                offset=request.offset,
                sources=[{
                    "entity_type": sql_result.get('entity_type', 'sql_query'),
                    "result_count": len(data_list),
                    "total_count": total_count,
                    "sql_preview": sql_query[:200],
                }],
            )
            logger.info(f"💾 Saved to history. Total messages in session: {len(get_recent_session_log(current_user.id, session_id))}")
            
            # Build response with user-friendly source reference
            entity_label = (sql_result.get('entity_type') or 'data').replace('_', ' ').title()
            sources = [ChatSource(
                rank=1,
                entity_type=sql_result.get('entity_type', 'platform_data'),
                entity_id="",
                framework_code="Live Data",
                control_code=None,
                control_name=f"{entity_label} — {len(data_list)} item(s) found",
                relevance_score=1.0,
                snippet=f"{len(data_list)} record(s) retrieved from your platform"
            )] if request.include_sources else []
            
            return ChatResponse(
                answer=answer,
                sources=sources,
                framework_filter="SQL_DIRECT_QUERY",
                timestamp=datetime.utcnow().isoformat(),
                has_more=has_more,
                total_count=total_count,
                current_offset=request.offset
            )
            
        except Exception as sql_error:
            # CRITICAL: Rollback transaction on error to prevent InFailedSqlTransaction state
            try:
                db.rollback()
                logger.info("[REFRESH] Transaction rolled back after SQL error")
            except Exception as rollback_err:
                logger.warning(f"[WARN]️ Rollback warning: {rollback_err}")
            
            logger.error(f"[FAIL] SQL execution error: {sql_error}")
            
            # [REFRESH] SMART RETRY: Check if it's a "no such column" error
            error_str = str(sql_error).lower()
            if 'no such column' in error_str or 'no such function' in error_str:
                logger.info("[REFRESH] DETECTED SCHEMA ERROR - Attempting smart retry...")
                
                # Extract table names from query
                import re
                table_pattern = r'FROM\s+(\w+)|JOIN\s+(\w+)'
                tables = set()
                for match in re.finditer(table_pattern, sql_query, re.IGNORECASE):
                    table_name = match.group(1) or match.group(2)
                    if table_name:
                        tables.add(table_name)
                
                # Fetch real schemas
                updated_schema_info = []
                for table in tables:
                    try:
                        schema_text = fetch_table_schema_from_db(table)
                        if schema_text:
                            updated_schema_info.append(schema_text)
                            logger.info(f"  [YES] Fetched {table} schema")
                    except Exception as e:
                        logger.warning(f"  [WARN]️ Failed to fetch {table} schema: {e}")
                
                if updated_schema_info:
                    # Retry with correct schema
                    retry_prompt = f"""
PREVIOUS QUERY FAILED WITH ERROR: {sql_error}

ACTUAL TABLE SCHEMAS (from database):
{chr(10).join(updated_schema_info)}

ORIGINAL QUESTION: {request.message}

Generate new SQL using ONLY the column names listed above. Use PostgreSQL syntax (NOW(), CURRENT_DATE, INTERVAL, ILIKE, DATE_TRUNC, TO_CHAR).
"""
                    
                    with usage_scope(
                        module_key="complychat",
                        feature_key="sql_generation_retry",
                        actor_user_id=current_user.id,
                        actor_username=current_user.username,
                    ):
                        retry_result = generate_sql_query(
                            retry_prompt,
                            language="en",
                            offset=request.offset,
                            limit=request.limit
                        )
                    
                    if retry_result and retry_result.get('sql'):
                        retry_sql = retry_result['sql']
                        logger.info(f"[REFRESH] RETRY SQL: {retry_sql[:150]}...")
                        
                        try:
                            # Try executing retry query
                            retry_exec = db.execute(text(retry_sql))
                            retry_rows = retry_exec.fetchall()
                            retry_data = [dict(row._mapping) for row in retry_rows]
                            
                            logger.info(f"[YES] RETRY SUCCESSFUL: {len(retry_data)} rows")
                            
                            # Check if empty result
                            if len(retry_data) == 0:
                                logger.info("[SAFE-EMPTY] Retry returned 0 rows")
                                answer = _no_platform_data_answer(request.message)
                            else:
                                # Format successful retry results
                                with usage_scope(
                                    module_key="complychat",
                                    feature_key="sql_result_formatting",
                                    actor_user_id=current_user.id,
                                    actor_username=current_user.username,
                                ):
                                    answer = format_query_results(retry_data, request.message, retry_sql, language="en")
                            
                            return ChatResponse(
                                answer=answer,
                                sources=[],
                                framework_filter=request.framework,
                                timestamp=datetime.utcnow().isoformat(),
                                has_more=False,
                                total_count=len(retry_data),
                                current_offset=0
                            )
                        except Exception as retry_error:
                            logger.error(f"[FAIL] RETRY EXECUTION FAILED: {retry_error}")
                            db.rollback()
            
            # If retry failed or not a schema error, return friendly message
            logger.info("[REFRESH] No retry possible or retry failed - returning friendly error")
            
            # Return deterministic failure to avoid fabricated examples.
            logger.info(f"[SAFE-FAIL] SQL execution failed; returning deterministic response: {error_str[:80]}")
            answer = _sql_failure_answer(request.message)
            
            # Save error to history
            store_chat_exchange(
                current_user.id,
                session_id,
                request.message,
                answer[:500],
                offset=request.offset,
                is_error=True,
                sources=[{"sql_preview": sql_query, "auto_retry_attempted": True}],
            )
            
            return ChatResponse(
                answer=answer,
                sources=[],
                framework_filter=request.framework,
                timestamp=datetime.utcnow().isoformat(),
                has_more=False,
                total_count=0,
                current_offset=0
            )
        
    except Exception as e:
        logger.error(f"Error in complychat ask: {e}", exc_info=True)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to process question: {str(e)}"
        )


@router.get("/history/{session_id}")
async def get_conversation_history(
    session_id: str,
    current_user: GRCUser = Depends(require_auth)
):
    """
    Get conversation history for a session.

    Returns the recent conversation context that the AI uses to understand follow-up questions.
    """
    history = get_recent_session_log(current_user.id, session_id)
    uploaded_files = load_session_uploaded_files(current_user.id, session_id)
    return {
        "session_id": resolve_chat_session_id(session_id, current_user.id),
        "messages": history,
        "count": len(history),
        "max_length": MAX_HISTORY_LENGTH,
        "uploaded_files": uploaded_files,
    }


@router.delete("/history/{session_id}")
async def clear_conversation_history(
    session_id: str,
    current_user: GRCUser = Depends(require_auth)
):
    """
    Clear conversation history for a session.

    Use this to delete a chat and its associated uploaded file context.
    """
    storage_key = get_session_storage_key(current_user.id, session_id)
    session_message_logs.pop(storage_key, None)
    conversation_history.pop(storage_key, None)
    session_uploaded_files.pop(storage_key, None)
    langchain_histories.pop(storage_key, None)

    upload_dir = get_session_upload_dir(current_user.id, session_id)
    if upload_dir.exists():
        for child in upload_dir.iterdir():
            if child.is_file():
                child.unlink(missing_ok=True)
        try:
            upload_dir.rmdir()
        except OSError:
            pass

    logger.info(f"🗑️ Cleared history for session: {session_id}")
    return {"message": f"History cleared for session {session_id}", "success": True}


@router.post("/upload")
async def upload_chat_files(
    files: List[UploadFile] = File(...),
    session_id: str = Form(...),
    db: Session = Depends(get_db),
    current_user: GRCUser = Depends(require_auth)
):
    """Upload one or more files for ComplyChat context. Any format is accepted."""
    if not files:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="No files were provided."
        )

    resolved_session_id = resolve_chat_session_id(session_id, current_user.id)
    upload_dir = get_session_upload_dir(current_user.id, resolved_session_id)
    existing_files = load_session_uploaded_files(current_user.id, resolved_session_id)
    uploaded_items: List[Dict[str, Any]] = []

    for upload in files:
        file_bytes = await upload.read()
        original_name = upload.filename or f"upload-{uuid4().hex}"
        saved_name = f"{uuid4().hex}_{original_name}"
        file_path = upload_dir / saved_name
        file_path.write_bytes(file_bytes)

        extracted_excerpt = extract_text_from_upload(original_name, file_bytes)
        file_info = {
            "id": uuid4().hex,
            "filename": original_name,
            "saved_name": saved_name,
            "path": str(file_path),
            "size": len(file_bytes),
            "content_type": upload.content_type or "application/octet-stream",
            "uploaded_at": datetime.utcnow().isoformat(),
            "excerpt": extracted_excerpt,
        }
        existing_files.append(file_info)
        uploaded_items.append(file_info)

    save_session_uploaded_files(current_user.id, resolved_session_id, existing_files)

    qdrant_indexed_points = 0
    service = get_qdrant_service()
    if service and getattr(service, "is_available", False):
        try:
            tenant_ids = get_user_tenants(current_user, db)
            qdrant_indexed_points = _index_chat_uploads_to_qdrant(
                service=service,
                tenant_ids=tenant_ids,
                uploaded_items=uploaded_items,
            )
        except Exception as exc:
            logger.warning("Chat upload indexing failed: %s", exc)

    return {
        "session_id": resolved_session_id,
        "count": len(uploaded_items),
        "uploaded_files": uploaded_items,
        "qdrant_indexed_points": qdrant_indexed_points,
    }


@router.get("/frameworks", response_model=List[FrameworkInfo])
def list_available_frameworks(
    db: Session = Depends(get_db),
    current_user: GRCUser = Depends(require_auth)
):
    """
    List all compliance frameworks available in the knowledge base.
    
    Use these framework codes to filter questions to specific frameworks.
    """
    try:
        # Query frameworks directly from database
        result = db.execute(text("""
            SELECT id, short_code, name, version, description 
            FROM grc_frameworks 
            WHERE is_active = true 
            ORDER BY name
        """))
        frameworks_data = result.fetchall()
        
        frameworks = [
            FrameworkInfo(
                id=str(row.id),
                code=row.short_code,
                name=row.name,
                version=row.version or "",
                description=row.description or ""
            )
            for row in frameworks_data
        ]
        
        logger.info(f"[YES] Found {len(frameworks)} active frameworks")
        return frameworks
    except Exception as e:
        logger.error(f"Error getting frameworks: {e}", exc_info=True)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to retrieve frameworks: {str(e)}"
        )


@router.get("/stats", response_model=StatsResponse)
def get_knowledge_base_stats(
    db: Session = Depends(get_db),
    current_user: GRCUser = Depends(require_auth)
):
    """
    Get statistics about the GRC database.
    
    Returns counts of key entities and list of available frameworks.
    """
    try:
        # Query statistics from database
        stats_query = text("""
            SELECT 
                (SELECT COUNT(*) FROM grc_frameworks WHERE is_active = true) as frameworks_count,
                (SELECT COUNT(*) FROM grc_framework_controls) as controls_count,
                (SELECT COUNT(*) FROM grc_vulnerabilities) as vulnerabilities_count,
                (SELECT COUNT(*) FROM grc_it_assets) as assets_count,
                (SELECT COUNT(*) FROM grc_risks) as risks_count
        """)
        stats_result = db.execute(stats_query)
        stats_row = stats_result.fetchone()
        
        # Query frameworks
        frameworks_query = text("""
            SELECT id, short_code, name, version, description 
            FROM grc_frameworks 
            WHERE is_active = true 
            ORDER BY name
        """)
        frameworks_result = db.execute(frameworks_query)
        frameworks_data = frameworks_result.fetchall()
        
        frameworks = [
            FrameworkInfo(
                id=str(row.id),
                code=row.short_code,
                name=row.name,
                version=row.version or "",
                description=row.description or ""
            )
            for row in frameworks_data
        ]
        
        total = sum([stats_row.frameworks_count, stats_row.controls_count, 
                     stats_row.vulnerabilities_count, stats_row.assets_count, stats_row.risks_count])
        
        return StatsResponse(
            total_entities=total,
            entity_types={
                "frameworks": stats_row.frameworks_count,
                "controls": stats_row.controls_count,
                "vulnerabilities": stats_row.vulnerabilities_count,
                "assets": stats_row.assets_count,
                "risks": stats_row.risks_count
            },
            available_frameworks=frameworks
        )
    except Exception as e:
        logger.error(f"Error getting stats: {e}", exc_info=True)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to retrieve statistics: {str(e)}"
        )


@router.post("/trigger-embedding-update")
def trigger_embedding_update(
    request: TriggerEmbeddingRequest,
    db: Session = Depends(get_db),
    current_user: GRCUser = Depends(require_auth)
):
    """
    Force-sync vector embeddings to Qdrant for the authenticated tenant(s).
    """
    service = get_qdrant_service()
    if not service or not getattr(service, "is_available", False):
        return {
            "status": "unavailable",
            "message": "Qdrant/OpenAI vector service is not configured.",
            "indexed_points": 0,
            "indexed_documents": 0,
        }

    tenant_ids = get_user_tenants(current_user, db)
    if not tenant_ids:
        return {
            "status": "ok",
            "message": "No tenant scope found for current user.",
            "indexed_points": 0,
            "indexed_documents": 0,
        }

    catalog = _build_vector_catalog(db, tenant_ids, uploaded_files=[])
    requested_types = {item.strip().lower() for item in (request.entity_types or []) if item and item.strip()}
    if requested_types:
        expanded_requested_types = set(requested_types)
        if "seeded_framework_control" in requested_types:
            expanded_requested_types.add("seeded_framework")
        if "parsed_framework_control" in requested_types:
            expanded_requested_types.add("uploaded_framework")
        if "policy_statement_version" in requested_types:
            expanded_requested_types.add("policy_statement")
        if "governance_document_version" in requested_types:
            expanded_requested_types.add("governance_document")
        selected = [
            entry for entry in catalog
            if str(entry.get("source_type", "")).strip().lower() in expanded_requested_types
        ]
    else:
        selected = catalog

    documents = _collect_documents_for_matches(
        db=db,
        tenant_ids=tenant_ids,
        matches=selected,
        uploaded_files=[],
    )
    indexed_points = _sync_docs_to_qdrant(service, documents, force=True)

    return {
        "status": "ok",
        "message": "Qdrant vector sync completed.",
        "indexed_points": indexed_points,
        "indexed_documents": len(documents),
        "selected_source_types": sorted({entry.get("source_type") for entry in selected}),
    }


# @router.post("/regenerate", response_model=dict)
# def regenerate_embeddings(
#     request: RegenerateRequest,
#     db: Session = Depends(get_db),
#     current_user: GRCUser = Depends(require_auth)
# ):
#     """OLD ENDPOINT - Kept for backward compatibility but returns deprecation notice"""
#     pass


@router.get("/health")
def health_check():
    """Check if SQL Agent service is healthy and ready"""
    qdrant_state = {"available": False, "reason": "not_initialized"}
    service = get_qdrant_service()
    if service:
        try:
            qdrant_state = service.health()
        except Exception as exc:
            qdrant_state = {"available": False, "reason": str(exc)}

    if not SQL_AGENT_ENABLED:
        return {
            "status": "unavailable",
            "message": "SQL Agent not initialized.",
            "ready": False,
            "vector": qdrant_state,
        }
    
    return {
        "status": "healthy",
        "message": "SQL Agent service is operational with optional Qdrant vector retrieval",
        "ready": True,
        "mode": "hybrid_sql_qdrant" if qdrant_state.get("available") else "pure_sql",
        "vector": qdrant_state,
    }


